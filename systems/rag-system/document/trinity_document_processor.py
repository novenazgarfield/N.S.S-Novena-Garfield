"""
🔄 三位一体文档处理器 (Trinity Document Processor)
================================================

专为"中央情报大脑"设计的新一代文档处理系统
- 支持多种文档格式的智能解析
- 保持文档结构的完整性
- 为三位一体分块做准备

Author: N.S.S-Novena-Garfield Project
Version: 2.0.0 - "Genesis"
"""

import io
import fitz  # PyMuPDF
from pathlib import Path
from typing import List, Union, BinaryIO, Dict, Any
from docx import Document
import pandas as pd
from pptx import Presentation
from bs4 import BeautifulSoup
import hashlib
from datetime import datetime

from config import DocumentConfig
from utils.logger import logger

class DocumentMetadata:
    """文档元数据类"""
    
    def __init__(self, filename: str, file_path: str = None):
        self.filename = filename
        self.file_path = file_path
        self.file_extension = Path(filename).suffix.lower()
        self.processed_at = datetime.now().isoformat()
        self.file_size = None
        self.file_hash = None
        self.page_count = None
        self.word_count = None
        self.character_count = None
        
    def calculate_file_stats(self, content: str, file_data: bytes = None):
        """计算文件统计信息"""
        self.character_count = len(content)
        self.word_count = len(content.split())
        
        if file_data:
            self.file_size = len(file_data)
            self.file_hash = hashlib.md5(file_data).hexdigest()
    
    def to_dict(self) -> Dict[str, Any]:
        """转换为字典格式"""
        return {
            "filename": self.filename,
            "file_path": self.file_path,
            "file_extension": self.file_extension,
            "processed_at": self.processed_at,
            "file_size": self.file_size,
            "file_hash": self.file_hash,
            "page_count": self.page_count,
            "word_count": self.word_count,
            "character_count": self.character_count
        }

class TrinityDocumentProcessor:
    """三位一体文档处理器"""
    
    def __init__(self):
        self.supported_extensions = DocumentConfig.SUPPORTED_EXTENSIONS
        logger.info("🔄 三位一体文档处理器已初始化")
    
    def _ensure_readable(self, file) -> BinaryIO:
        """将文件统一为可被处理的 BytesIO 或文件路径"""
        try:
            if isinstance(file, (str, Path)):  # 本地路径
                return open(file, 'rb')
            elif hasattr(file, "getvalue"):  # Streamlit UploadedFile
                return io.BytesIO(file.getvalue())
            elif hasattr(file, "read"):  # 可能是 open() 后的文件句柄
                content = file.read()
                if isinstance(content, str):
                    content = content.encode('utf-8')
                return io.BytesIO(content)
            else:
                raise TypeError("无法识别的文件类型，请检查 file 对象。")
        except Exception as e:
            logger.error(f"文件转换失败: {e}")
            raise
    
    def _extract_pdf_with_structure(self, file) -> Dict[str, Any]:
        """提取PDF文件内容，保持结构信息"""
        try:
            file_like = self._ensure_readable(file)
            doc = fitz.open(stream=file_like, filetype="pdf")
            
            full_text = ""
            pages_content = []
            
            for page_num in range(len(doc)):
                page = doc[page_num]
                page_text = page.get_text("text")
                pages_content.append({
                    "page_number": page_num + 1,
                    "content": page_text
                })
                full_text += f"\n\n--- 第 {page_num + 1} 页 ---\n\n{page_text}"
            
            doc.close()
            
            result = {
                "content": full_text.strip(),
                "structure": {
                    "type": "pdf",
                    "page_count": len(pages_content),
                    "pages": pages_content
                }
            }
            
            logger.info(f"PDF文件解析成功: {len(pages_content)} 页")
            return result
            
        except Exception as e:
            logger.error(f"PDF文件解析失败: {e}")
            raise
    
    def _extract_docx_with_structure(self, file) -> Dict[str, Any]:
        """提取DOCX文件内容，保持段落结构"""
        try:
            file_like = self._ensure_readable(file)
            doc = Document(file_like)
            
            paragraphs_content = []
            full_text = ""
            
            for para_idx, para in enumerate(doc.paragraphs):
                if para.text.strip():  # 只处理非空段落
                    paragraphs_content.append({
                        "paragraph_number": para_idx + 1,
                        "content": para.text.strip(),
                        "style": para.style.name if para.style else "Normal"
                    })
                    full_text += f"{para.text.strip()}\n\n"
            
            result = {
                "content": full_text.strip(),
                "structure": {
                    "type": "docx",
                    "paragraph_count": len(paragraphs_content),
                    "paragraphs": paragraphs_content
                }
            }
            
            logger.info(f"DOCX文件解析成功: {len(paragraphs_content)} 个段落")
            return result
            
        except Exception as e:
            logger.error(f"DOCX文件解析失败: {e}")
            raise
    
    def _extract_pptx_with_structure(self, file) -> Dict[str, Any]:
        """提取PPTX文件内容，保持幻灯片结构"""
        try:
            file_like = self._ensure_readable(file)
            prs = Presentation(file_like)
            
            slides_content = []
            full_text = ""
            
            for slide_idx, slide in enumerate(prs.slides):
                slide_text = ""
                shapes_content = []
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        shapes_content.append(shape.text.strip())
                        slide_text += f"{shape.text.strip()}\n"
                
                if slide_text.strip():
                    slides_content.append({
                        "slide_number": slide_idx + 1,
                        "content": slide_text.strip(),
                        "shapes_count": len(shapes_content)
                    })
                    full_text += f"\n\n--- 幻灯片 {slide_idx + 1} ---\n\n{slide_text.strip()}"
            
            result = {
                "content": full_text.strip(),
                "structure": {
                    "type": "pptx",
                    "slide_count": len(slides_content),
                    "slides": slides_content
                }
            }
            
            logger.info(f"PPTX文件解析成功: {len(slides_content)} 张幻灯片")
            return result
            
        except Exception as e:
            logger.error(f"PPTX文件解析失败: {e}")
            raise
    
    def _extract_excel_with_structure(self, file) -> Dict[str, Any]:
        """提取Excel文件内容，保持工作表结构"""
        try:
            file_like = self._ensure_readable(file)
            
            # 读取所有工作表
            excel_file = pd.ExcelFile(file_like)
            sheets_content = []
            full_text = ""
            
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(file_like, sheet_name=sheet_name)
                sheet_text = df.to_string(index=False)
                
                sheets_content.append({
                    "sheet_name": sheet_name,
                    "content": sheet_text,
                    "rows": len(df),
                    "columns": len(df.columns)
                })
                
                full_text += f"\n\n--- 工作表: {sheet_name} ---\n\n{sheet_text}"
            
            result = {
                "content": full_text.strip(),
                "structure": {
                    "type": "excel",
                    "sheet_count": len(sheets_content),
                    "sheets": sheets_content
                }
            }
            
            logger.info(f"Excel文件解析成功: {len(sheets_content)} 个工作表")
            return result
            
        except Exception as e:
            logger.error(f"Excel文件解析失败: {e}")
            raise
    
    def _extract_csv_with_structure(self, file) -> Dict[str, Any]:
        """提取CSV文件内容"""
        try:
            file_like = self._ensure_readable(file)
            df = pd.read_csv(file_like)
            content = df.to_string(index=False)
            
            result = {
                "content": content,
                "structure": {
                    "type": "csv",
                    "rows": len(df),
                    "columns": len(df.columns),
                    "column_names": df.columns.tolist()
                }
            }
            
            logger.info(f"CSV文件解析成功: {len(df)} 行 x {len(df.columns)} 列")
            return result
            
        except Exception as e:
            logger.error(f"CSV文件解析失败: {e}")
            raise
    
    def _extract_text_with_structure(self, file) -> Dict[str, Any]:
        """提取文本文件内容，保持行结构"""
        try:
            file_like = self._ensure_readable(file)
            content = file_like.read().decode("utf-8", errors="ignore")
            
            lines = content.split('\n')
            non_empty_lines = [line for line in lines if line.strip()]
            
            result = {
                "content": content,
                "structure": {
                    "type": "text",
                    "total_lines": len(lines),
                    "non_empty_lines": len(non_empty_lines)
                }
            }
            
            logger.info(f"文本文件解析成功: {len(lines)} 行")
            return result
            
        except Exception as e:
            logger.error(f"文本文件解析失败: {e}")
            raise
    
    def _extract_html_with_structure(self, file) -> Dict[str, Any]:
        """提取HTML文件内容，保持标签结构信息"""
        try:
            file_like = self._ensure_readable(file)
            soup = BeautifulSoup(file_like, "html.parser")
            
            # 提取纯文本
            content = soup.get_text()
            
            # 分析HTML结构
            structure_info = {
                "type": "html",
                "title": soup.title.string if soup.title else None,
                "headings": {
                    "h1": len(soup.find_all('h1')),
                    "h2": len(soup.find_all('h2')),
                    "h3": len(soup.find_all('h3')),
                    "h4": len(soup.find_all('h4')),
                    "h5": len(soup.find_all('h5')),
                    "h6": len(soup.find_all('h6'))
                },
                "paragraphs": len(soup.find_all('p')),
                "links": len(soup.find_all('a')),
                "images": len(soup.find_all('img'))
            }
            
            result = {
                "content": content,
                "structure": structure_info
            }
            
            logger.info("HTML文件解析成功")
            return result
            
        except Exception as e:
            logger.error(f"HTML文件解析失败: {e}")
            raise
    
    def process_document(self, file, filename: str = None) -> Dict[str, Any]:
        """处理单个文档，返回结构化结果"""
        try:
            # 确定文件名
            if filename is None:
                if hasattr(file, 'name'):
                    filename = file.name
                elif isinstance(file, (str, Path)):
                    filename = Path(file).name
                else:
                    filename = "unknown_document"
            
            # 获取文件扩展名
            ext = Path(filename).suffix.lower()
            
            logger.info(f"开始处理文档: {filename} (类型: {ext})")
            
            # 检查是否支持该文件类型
            if ext not in self.supported_extensions:
                raise ValueError(f"不支持的文件类型：{ext}")
            
            # 创建文档元数据
            metadata = DocumentMetadata(filename)
            
            # 根据文件类型调用相应的解析器
            if ext == ".pdf":
                extraction_result = self._extract_pdf_with_structure(file)
            elif ext == ".docx":
                extraction_result = self._extract_docx_with_structure(file)
            elif ext == ".pptx":
                extraction_result = self._extract_pptx_with_structure(file)
            elif ext in [".xlsx", ".xls"]:
                extraction_result = self._extract_excel_with_structure(file)
            elif ext == ".csv":
                extraction_result = self._extract_csv_with_structure(file)
            elif ext in [".txt", ".py", ".md"]:
                extraction_result = self._extract_text_with_structure(file)
            elif ext == ".html":
                extraction_result = self._extract_html_with_structure(file)
            elif ext == ".caj":
                raise NotImplementedError("暂不支持 caj 文件，请先转换为 PDF。")
            else:
                raise ValueError(f"不支持的文件类型：{ext}")
            
            # 计算文件统计信息
            content = extraction_result["content"]
            if hasattr(file, 'getvalue'):
                file_data = file.getvalue()
            elif isinstance(file, (str, Path)):
                with open(file, 'rb') as f:
                    file_data = f.read()
            else:
                file_data = None
            
            metadata.calculate_file_stats(content, file_data)
            
            # 如果是PDF，更新页数信息
            if ext == ".pdf" and "structure" in extraction_result:
                metadata.page_count = extraction_result["structure"].get("page_count", 0)
            
            # 构建最终结果
            result = {
                "success": True,
                "filename": filename,
                "content": content,
                "metadata": metadata.to_dict(),
                "structure": extraction_result.get("structure", {}),
                "processing_time": datetime.now().isoformat()
            }
            
            logger.info(f"文档处理成功: {filename} ({len(content)} 字符)")
            return result
            
        except Exception as e:
            logger.error(f"文档处理失败 {filename}: {e}")
            return {
                "success": False,
                "filename": filename,
                "error": str(e),
                "processing_time": datetime.now().isoformat()
            }
    
    def process_documents_batch(self, files: List) -> Dict[str, Any]:
        """批量处理文档"""
        try:
            logger.info(f"开始批量处理 {len(files)} 个文档")
            
            successful_results = []
            failed_results = []
            total_characters = 0
            
            for file in files:
                result = self.process_document(file)
                
                if result["success"]:
                    successful_results.append(result)
                    total_characters += len(result["content"])
                else:
                    failed_results.append(result)
            
            batch_result = {
                "success": len(successful_results) > 0,
                "total_files": len(files),
                "successful_count": len(successful_results),
                "failed_count": len(failed_results),
                "total_characters": total_characters,
                "successful_results": successful_results,
                "failed_results": failed_results,
                "batch_processing_time": datetime.now().isoformat()
            }
            
            logger.info(f"批量处理完成: 成功 {len(successful_results)}, 失败 {len(failed_results)}")
            
            if failed_results:
                failed_files = [r["filename"] for r in failed_results]
                logger.warning(f"处理失败的文件: {failed_files}")
            
            return batch_result
            
        except Exception as e:
            logger.error(f"批量文档处理失败: {e}")
            return {
                "success": False,
                "error": str(e),
                "batch_processing_time": datetime.now().isoformat()
            }