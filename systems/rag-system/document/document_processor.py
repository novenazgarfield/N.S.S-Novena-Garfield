"""
文档处理系统
"""
import io
from pathlib import Path
from typing import List, Union, BinaryIO
import fitz  # PyMuPDF
from docx import Document
import pandas as pd
from pptx import Presentation
from bs4 import BeautifulSoup

from config import DocumentConfig
from utils.logger import logger

class DocumentProcessor:
    """文档处理类"""
    
    def __init__(self):
        self.supported_extensions = DocumentConfig.SUPPORTED_EXTENSIONS
    
    def ensure_readable(self, file) -> BinaryIO:
        """将文件统一为可被处理的 BytesIO 或文件路径"""
        try:
            if isinstance(file, (str, Path)):  # 本地路径
                return open(file, 'rb')
            elif hasattr(file, "getvalue"):  # Streamlit UploadedFile
                return io.BytesIO(file.getvalue())
            elif hasattr(file, "read"):  # 可能是 open() 后的文件句柄
                content = file.read()
                if isinstance(content, str):
                    content = content.encode('utf-8')
                return io.BytesIO(content)
            else:
                raise TypeError("无法识别的文件类型，请检查 file 对象。")
        except Exception as e:
            logger.error(f"文件转换失败: {e}")
            raise
    
    def extract_pdf(self, file) -> str:
        """提取PDF文件内容"""
        try:
            file_like = self.ensure_readable(file)
            doc = fitz.open(stream=file_like, filetype="pdf")
            text = ""
            for page in doc:
                text += page.get_text("text")
            doc.close()
            logger.info("PDF文件解析成功")
            return text
        except Exception as e:
            logger.error(f"PDF文件解析失败: {e}")
            raise
    
    def extract_docx(self, file) -> str:
        """提取DOCX文件内容"""
        try:
            file_like = self.ensure_readable(file)
            doc = Document(file_like)
            text = "\n".join([para.text for para in doc.paragraphs])
            logger.info("DOCX文件解析成功")
            return text
        except Exception as e:
            logger.error(f"DOCX文件解析失败: {e}")
            raise
    
    def extract_pptx(self, file) -> str:
        """提取PPTX文件内容"""
        try:
            file_like = self.ensure_readable(file)
            prs = Presentation(file_like)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            logger.info("PPTX文件解析成功")
            return text
        except Exception as e:
            logger.error(f"PPTX文件解析失败: {e}")
            raise
    
    def extract_excel(self, file) -> str:
        """提取Excel文件内容"""
        try:
            file_like = self.ensure_readable(file)
            df = pd.read_excel(file_like)
            text = df.to_string(index=False)
            logger.info("Excel文件解析成功")
            return text
        except Exception as e:
            logger.error(f"Excel文件解析失败: {e}")
            raise
    
    def extract_csv(self, file) -> str:
        """提取CSV文件内容"""
        try:
            file_like = self.ensure_readable(file)
            df = pd.read_csv(file_like)
            text = df.to_string(index=False)
            logger.info("CSV文件解析成功")
            return text
        except Exception as e:
            logger.error(f"CSV文件解析失败: {e}")
            raise
    
    def extract_text(self, file) -> str:
        """提取文本文件内容"""
        try:
            file_like = self.ensure_readable(file)
            text = file_like.read().decode("utf-8", errors="ignore")
            logger.info("文本文件解析成功")
            return text
        except Exception as e:
            logger.error(f"文本文件解析失败: {e}")
            raise
    
    def extract_html(self, file) -> str:
        """提取HTML文件内容"""
        try:
            file_like = self.ensure_readable(file)
            soup = BeautifulSoup(file_like, "html.parser")
            text = soup.get_text()
            logger.info("HTML文件解析成功")
            return text
        except Exception as e:
            logger.error(f"HTML文件解析失败: {e}")
            raise
    
    def read_file(self, file) -> str:
        """根据文件类型调用相应的解析器"""
        try:
            # 获取文件扩展名
            if hasattr(file, 'name'):
                ext = Path(file.name).suffix.lower()
                filename = file.name
            elif isinstance(file, (str, Path)):
                ext = Path(file).suffix.lower()
                filename = Path(file).name
            else:
                raise ValueError("无法确定文件类型")
            
            logger.info(f"正在处理文件: {filename}, 类型: {ext}")
            
            # 检查是否支持该文件类型
            if ext not in self.supported_extensions:
                raise ValueError(f"不支持的文件类型：{ext}")
            
            # 根据文件类型调用相应的解析器
            if ext == ".pdf":
                return self.extract_pdf(file)
            elif ext == ".docx":
                return self.extract_docx(file)
            elif ext == ".pptx":
                return self.extract_pptx(file)
            elif ext in [".xlsx", ".xls"]:
                return self.extract_excel(file)
            elif ext == ".csv":
                return self.extract_csv(file)
            elif ext in [".txt", ".py", ".md"]:
                return self.extract_text(file)
            elif ext == ".html":
                return self.extract_html(file)
            elif ext == ".caj":
                raise NotImplementedError("暂不支持 caj 文件，请先转换为 PDF。")
            else:
                raise ValueError(f"不支持的文件类型：{ext}")
                
        except Exception as e:
            logger.error(f"文件处理失败: {e}")
            raise
    
    def chunk_text(self, text: str, chunk_size: int = None, overlap: int = None) -> List[str]:
        """文本分块"""
        try:
            chunk_size = chunk_size or DocumentConfig.CHUNK_SIZE
            overlap = overlap or DocumentConfig.CHUNK_OVERLAP
            
            words = text.split()
            chunks = []
            i = 0
            while i < len(words):
                chunk = words[i:i + chunk_size]
                chunks.append(" ".join(chunk))
                i += chunk_size - overlap
            
            logger.info(f"文本分块完成: {len(chunks)} 个块")
            return chunks
            
        except Exception as e:
            logger.error(f"文本分块失败: {e}")
            raise
    
    def process_files(self, files: List) -> List[str]:
        """批量处理文件"""
        try:
            all_texts = []
            successful_files = []
            failed_files = []
            
            for file in files:
                try:
                    text = self.read_file(file)
                    if text.strip():  # 确保文本不为空
                        all_texts.append(text)
                        filename = getattr(file, 'name', str(file))
                        successful_files.append(filename)
                    else:
                        logger.warning(f"文件内容为空: {getattr(file, 'name', str(file))}")
                except Exception as e:
                    filename = getattr(file, 'name', str(file))
                    failed_files.append(filename)
                    logger.error(f"处理文件失败 {filename}: {e}")
            
            logger.info(f"文件处理完成: 成功 {len(successful_files)}, 失败 {len(failed_files)}")
            
            if failed_files:
                logger.warning(f"失败的文件: {failed_files}")
            
            return all_texts
            
        except Exception as e:
            logger.error(f"批量文件处理失败: {e}")
            raise