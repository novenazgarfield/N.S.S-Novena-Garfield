"""
RAG系统通用组件
包含所有平台共用的基础功能
"""

import streamlit as st
import time
import datetime
from datetime import timezone, timedelta
import hashlib
import os
from typing import Dict, List, Any, Optional

# 中国时区 (UTC+8)
CHINA_TZ = timezone(timedelta(hours=8))

def get_china_time():
    """获取中国时区的当前时间"""
    return datetime.datetime.now(CHINA_TZ)

# 导入可选依赖
try:
    import PyPDF2
    import docx
    import pandas as pd
    from pptx import Presentation
    from bs4 import BeautifulSoup
    FULL_FEATURES = True
except ImportError:
    FULL_FEATURES = False

class SystemConfig:
    """系统配置管理"""
    
    DEFAULT_CONFIG = {
        "system_name": "RAG智能对话",
        "version": "v3.2-Universal",
        "admin_password": "rag2024",
        "max_file_size": 50,  # MB
        "max_daily_queries": 500,
        "max_message_length": 5000,
        "max_documents": 20,
        "max_conversations": 1000,
        "chat_speed": 1.0,
        "theme_mode": "light",  # light, dark, auto
        "show_timestamps": True,
        "auto_scroll": True,
        "supported_formats": [".pdf", ".docx", ".txt", ".md", ".pptx", ".csv"]
    }
    
    def __init__(self):
        self.config = self.DEFAULT_CONFIG.copy()
    
    def get(self, key: str, default=None):
        return self.config.get(key, default)
    
    def set(self, key: str, value: Any):
        self.config[key] = value
    
    def update(self, updates: Dict[str, Any]):
        self.config.update(updates)
    
    def reset(self):
        self.config = self.DEFAULT_CONFIG.copy()

class DocumentProcessor:
    """文档处理器"""
    
    def __init__(self):
        self.supported_formats = [".pdf", ".docx", ".txt", ".md", ".pptx", ".csv"]
    
    def extract_text_from_file(self, uploaded_file) -> Optional[str]:
        """从上传的文件中提取文本"""
        try:
            file_extension = os.path.splitext(uploaded_file.name)[1].lower()
            
            if file_extension == ".txt" or file_extension == ".md":
                return str(uploaded_file.read(), "utf-8")
            
            elif file_extension == ".pdf" and FULL_FEATURES:
                return self._extract_from_pdf(uploaded_file)
            
            elif file_extension == ".docx" and FULL_FEATURES:
                return self._extract_from_docx(uploaded_file)
            
            elif file_extension == ".pptx" and FULL_FEATURES:
                return self._extract_from_pptx(uploaded_file)
            
            elif file_extension == ".csv" and FULL_FEATURES:
                return self._extract_from_csv(uploaded_file)
            
            else:
                return f"文件格式 {file_extension} 暂不支持，请上传 {', '.join(self.supported_formats)} 格式的文件。"
        
        except Exception as e:
            return f"文件处理出错: {str(e)}"
    
    def _extract_from_pdf(self, uploaded_file) -> str:
        """从PDF提取文本"""
        try:
            pdf_reader = PyPDF2.PdfReader(uploaded_file)
            text = ""
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
            return text.strip()
        except Exception as e:
            return f"PDF处理失败: {str(e)}"
    
    def _extract_from_docx(self, uploaded_file) -> str:
        """从Word文档提取文本"""
        try:
            doc = docx.Document(uploaded_file)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            return text.strip()
        except Exception as e:
            return f"Word文档处理失败: {str(e)}"
    
    def _extract_from_pptx(self, uploaded_file) -> str:
        """从PowerPoint提取文本"""
        try:
            prs = Presentation(uploaded_file)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text.strip()
        except Exception as e:
            return f"PowerPoint处理失败: {str(e)}"
    
    def _extract_from_csv(self, uploaded_file) -> str:
        """从CSV提取文本"""
        try:
            df = pd.read_csv(uploaded_file)
            return df.to_string()
        except Exception as e:
            return f"CSV处理失败: {str(e)}"

class ResponseGenerator:
    """AI回答生成器 - 集成多种AI模型"""
    
    def __init__(self):
        self.predefined_responses = {
            "你好": "您好！很高兴与您对话。我是RAG智能助手，可以帮您解答问题和分析文档。",
            "什么是RAG": "RAG（检索增强生成）是一种AI技术，结合了信息检索和文本生成。它先从知识库中检索相关信息，然后基于检索结果生成准确的回答。",
            "如何使用": "使用很简单！您可以：\n1. 💬 直接输入问题进行对话\n2. 📄 上传文档进行分析\n3. ⚙️ 点击设置按钮调整参数\n4. 🤖 选择不同的AI模型",
            "分析文档": "请先上传文档，我就可以为您分析文档内容了。支持PDF、Word、PPT等多种格式。",
            "系统设置": "点击右上角的⚙️设置按钮，可以调整AI速度、主题、文件限制等各种参数。",
            "AI模型": "系统支持多种AI模型：OpenAI GPT、Google Gemini、Anthropic Claude等。您可以在设置中配置API密钥并选择模型。"
        }
        
        # 尝试导入AI模型管理器
        try:
            from ai_models import get_ai_model_manager
            self.ai_manager = get_ai_model_manager()
            self.use_ai_models = True
        except ImportError:
            self.ai_manager = None
            self.use_ai_models = False
    
    def generate_response(self, prompt: str, context: str = "", user_id: str = None, **kwargs) -> str:
        """生成AI回答"""
        # 如果有AI模型管理器且用户已选择模型，使用AI生成
        if self.use_ai_models and self.ai_manager and self.ai_manager.current_model:
            try:
                return self.ai_manager.generate_response(prompt, context, user_id, **kwargs)
            except Exception as e:
                # AI生成失败时回退到预定义回答
                st.warning(f"AI模型调用失败，使用演示模式: {str(e)}")
        
        # 检查预定义回答
        for key, value in self.predefined_responses.items():
            if key in prompt:
                return value
        
        # 基于文档上下文的回答
        if context:
            return f"基于您上传的文档，我分析如下：\n\n{context[:500]}...\n\n这是文档的主要内容概述。您想了解文档的哪个具体方面？"
        
        # 默认智能回答
        return f"感谢您的问题：「{prompt}」\n\n这是一个很有意思的话题。基于RAG技术，我正在为您检索相关信息并生成回答。\n\n💡 **提示：** \n• 您可以上传相关文档，我能提供更准确的分析\n• 在设置中配置API密钥可使用真实的AI模型\n• 支持OpenAI、Gemini、Claude等多种模型"

class SessionManager:
    """会话状态管理"""
    
    @staticmethod
    def init_session_state():
        """初始化会话状态"""
        if "messages" not in st.session_state:
            st.session_state.messages = [
                {
                    "role": "assistant",
                    "content": "👋 您好！我是RAG智能对话助手。\n\n✨ **我可以帮您：**\n• 💬 智能问答对话\n• 📄 分析上传的文档\n• 🔍 检索相关信息\n\n请开始我们的对话吧！",
                    "timestamp": get_china_time().strftime("%H:%M")
                }
            ]
        
        if "uploaded_documents" not in st.session_state:
            st.session_state.uploaded_documents = {}
        
        if "query_count" not in st.session_state:
            st.session_state.query_count = 0
        
        if "show_settings" not in st.session_state:
            st.session_state.show_settings = False
        
        if "device_type" not in st.session_state:
            st.session_state.device_type = "auto"  # auto, mobile, desktop
    
    @staticmethod
    def add_message(role: str, content: str):
        """添加消息到会话"""
        message = {
            "role": role,
            "content": content,
            "timestamp": get_china_time().strftime("%H:%M")
        }
        st.session_state.messages.append(message)
    
    @staticmethod
    def clear_messages():
        """清空消息"""
        st.session_state.messages = []
    
    @staticmethod
    def add_document(name: str, content: str, size: int) -> str:
        """添加文档"""
        doc_id = hashlib.md5(f"{name}{time.time()}".encode()).hexdigest()[:8]
        st.session_state.uploaded_documents[doc_id] = {
            "name": name,
            "content": content,
            "upload_time": get_china_time().strftime("%m-%d %H:%M"),
            "size": size
        }
        return doc_id
    
    @staticmethod
    def remove_document(doc_id: str):
        """删除文档"""
        if doc_id in st.session_state.uploaded_documents:
            del st.session_state.uploaded_documents[doc_id]

class AuthManager:
    """认证管理"""
    
    @staticmethod
    def check_password(password: str = "rag2024") -> bool:
        """检查密码"""
        if "password_correct" not in st.session_state:
            st.session_state.password_correct = False
        
        if st.session_state.password_correct:
            return True
        
        st.markdown("# 🔐 访问验证")
        st.markdown("### 请输入访问密码以继续使用系统")
        
        input_password = st.text_input("密码", type="password", key="password_input")
        
        col1, col2 = st.columns(2)
        with col1:
            if st.button("🚀 进入系统", type="primary", use_container_width=True):
                if input_password == password:
                    st.session_state.password_correct = True
                    st.success("✅ 验证成功！正在进入系统...")
                    time.sleep(1)
                    st.rerun()
                else:
                    st.error("❌ 密码错误，请重试")
        
        with col2:
            if st.button("💡 获取帮助", use_container_width=True):
                st.info("💡 **提示：** 默认密码是 `rag2024`")
        
        st.markdown("---")
        st.markdown(
            "<div style='text-align: center; color: #666; font-size: 0.9rem;'>"
            "🔒 系统采用密码保护，确保安全访问<br>"
            "📱 支持移动端和桌面端全平台使用"
            "</div>",
            unsafe_allow_html=True
        )
        
        return False

class UIComponents:
    """UI组件库"""
    
    @staticmethod
    def render_message_bubble(message: Dict[str, Any], show_timestamp: bool = True):
        """渲染消息气泡"""
        timestamp = f" {message.get('timestamp', '')}" if show_timestamp else ""
        
        if message["role"] == "user":
            st.markdown(f"""
            <div style="background: linear-gradient(135deg, #e3f2fd 0%, #bbdefb 100%); 
                        padding: 12px 16px; border-radius: 18px 18px 4px 18px; 
                        margin: 8px 0; margin-left: 15%; box-shadow: 0 2px 8px rgba(0,0,0,0.1);">
                <div style="font-weight: 600; color: #1976d2; font-size: 0.9rem;">
                    👤 您{timestamp}
                </div>
                <div style="margin-top: 4px; line-height: 1.5;">
                    {message['content']}
                </div>
            </div>
            """, unsafe_allow_html=True)
        else:
            st.markdown(f"""
            <div style="background: linear-gradient(135deg, #f5f5f5 0%, #eeeeee 100%); 
                        padding: 12px 16px; border-radius: 18px 18px 18px 4px; 
                        margin: 8px 0; margin-right: 15%; box-shadow: 0 2px 8px rgba(0,0,0,0.1);">
                <div style="font-weight: 600; color: #388e3c; font-size: 0.9rem;">
                    🤖 AI助手{timestamp}
                </div>
                <div style="margin-top: 4px; line-height: 1.5;">
                    {message['content']}
                </div>
            </div>
            """, unsafe_allow_html=True)
    
    @staticmethod
    def render_document_card(doc_id: str, doc_info: Dict[str, Any], show_delete: bool = True):
        """渲染文档卡片"""
        col1, col2 = st.columns([3, 1] if show_delete else [1])
        
        with col1:
            file_size_mb = doc_info.get('size', 0) / (1024 * 1024)
            st.markdown(f"""
            <div style="background: #f8f9fa; padding: 10px; border-radius: 8px; border-left: 4px solid #007bff;">
                <div style="font-weight: 600; color: #007bff;">
                    📄 {doc_info['name'][:40]}{'...' if len(doc_info['name']) > 40 else ''}
                </div>
                <div style="font-size: 0.8rem; color: #666; margin-top: 4px;">
                    📅 {doc_info['upload_time']} | 📊 {file_size_mb:.2f} MB
                </div>
            </div>
            """, unsafe_allow_html=True)
        
        if show_delete:
            with col2:
                if st.button("🗑️", key=f"del_{doc_id}", help="删除文档"):
                    SessionManager.remove_document(doc_id)
                    st.success("✅ 文档已删除")
                    st.rerun()
    
    @staticmethod
    def render_system_stats(config: SystemConfig):
        """渲染系统统计"""
        col1, col2, col3, col4 = st.columns(4)
        
        with col1:
            st.metric(
                "💬 对话轮次", 
                len(st.session_state.messages) // 2,
                help="当前会话的对话轮次"
            )
        
        with col2:
            st.metric(
                "📄 文档数量", 
                len(st.session_state.uploaded_documents),
                help="已上传的文档数量"
            )
        
        with col3:
            st.metric(
                "📈 查询次数", 
                st.session_state.query_count,
                help="今日查询次数"
            )
        
        with col4:
            usage_percent = (st.session_state.query_count / config.get("max_daily_queries")) * 100
            st.metric(
                "📊 使用率", 
                f"{usage_percent:.1f}%",
                help="今日查询使用率"
            )

# 全局实例
system_config = SystemConfig()
document_processor = DocumentProcessor()
response_generator = ResponseGenerator()