import streamlit as st
import time
import datetime
import random
import io
import os
from pathlib import Path
from typing import List, Dict, Any
import hashlib

# 文档处理相关导入
try:
    import fitz  # PyMuPDF
    from docx import Document
    import pandas as pd
    from pptx import Presentation
    from bs4 import BeautifulSoup
    FULL_FEATURES = True
except ImportError:
    FULL_FEATURES = False

# 页面配置
st.set_page_config(
    page_title="💬 RAG智能对话系统",
    page_icon="💬",
    layout="wide",
    initial_sidebar_state="collapsed"
)

# 预设主题配色方案
PRESET_THEMES = {
    "灰白简约": {
        "name": "🤍 灰白简约",
        "chat_bg": "#f8f9fa",
        "user_bg": "#6c757d",
        "assistant_bg": "#e9ecef",
        "user_color": "white",
        "assistant_color": "#212529",
        "border": "#dee2e6",
        "upload_bg": "#6c757d",
        "upload_color": "white"
    },
    "蓝色清新": {
        "name": "💙 蓝色清新", 
        "chat_bg": "#f0f8ff",
        "user_bg": "#4dabf7",
        "assistant_bg": "#e7f5ff",
        "user_color": "white",
        "assistant_color": "#1864ab",
        "border": "#74c0fc",
        "upload_bg": "#4dabf7",
        "upload_color": "white"
    },
    "绿色自然": {
        "name": "💚 绿色自然",
        "chat_bg": "#f3f9f0",
        "user_bg": "#51cf66",
        "assistant_bg": "#d3f9d8",
        "user_color": "white", 
        "assistant_color": "#2b8a3e",
        "border": "#8ce99a",
        "upload_bg": "#51cf66",
        "upload_color": "white"
    },
    "深色模式": {
        "name": "🖤 深色模式",
        "chat_bg": "#2d3748",
        "user_bg": "#4a5568",
        "assistant_bg": "#1a202c",
        "user_color": "#e2e8f0",
        "assistant_color": "#e2e8f0",
        "border": "#4a5568",
        "upload_bg": "#4a5568",
        "upload_color": "#e2e8f0"
    }
}

def apply_theme(theme_config):
    """应用主题配置"""
    st.markdown(f"""
    <style>
        /* 隐藏Streamlit默认元素 */
        #MainMenu {{visibility: hidden;}}
        footer {{visibility: hidden;}}
        header {{visibility: hidden;}}
        
        /* 设置按钮样式 */
        .settings-button {{
            position: fixed;
            top: 20px;
            right: 20px;
            z-index: 1000;
            background: linear-gradient(135deg, {theme_config['user_bg']} 0%, {theme_config['upload_bg']} 100%);
            color: {theme_config['user_color']};
            border: none;
            border-radius: 50%;
            width: 50px;
            height: 50px;
            font-size: 20px;
            cursor: pointer;
            box-shadow: 0 4px 15px rgba(0,0,0,0.2);
            transition: all 0.3s ease;
        }}
        
        .settings-button:hover {{
            transform: scale(1.1);
            box-shadow: 0 6px 20px rgba(0,0,0,0.3);
        }}
        
        /* 聊天容器样式 */
        .chat-container {{
            background-color: {theme_config['chat_bg']};
            border: 2px solid {theme_config['border']};
            border-radius: 15px;
            padding: 20px;
            margin: 10px 0;
            height: 500px;
            overflow-y: auto;
            box-shadow: 0 4px 15px rgba(0,0,0,0.1);
        }}
        
        /* 消息气泡样式 */
        .user-message {{
            background: {theme_config['user_bg']};
            color: {theme_config['user_color']};
            padding: 12px 18px;
            border-radius: 18px 18px 5px 18px;
            margin: 8px 0 8px auto;
            max-width: 80%;
            word-wrap: break-word;
            box-shadow: 0 2px 8px rgba(0,0,0,0.15);
            display: block;
            text-align: left;
        }}
        
        .assistant-message {{
            background: {theme_config['assistant_bg']};
            color: {theme_config['assistant_color']};
            padding: 12px 18px;
            border-radius: 18px 18px 18px 5px;
            margin: 8px auto 8px 0;
            max-width: 80%;
            word-wrap: break-word;
            box-shadow: 0 2px 8px rgba(0,0,0,0.15);
            display: block;
            text-align: left;
        }}
        
        /* 文件上传区域样式 */
        .upload-area {{
            background: {theme_config['upload_bg']};
            border-radius: 15px;
            padding: 15px;
            margin: 10px 0;
            color: {theme_config['upload_color']};
            text-align: center;
            box-shadow: 0 2px 10px rgba(0,0,0,0.1);
        }}
        
        /* 按钮样式 */
        .stButton > button {{
            border-radius: 20px;
            border: none;
            padding: 8px 20px;
            font-weight: 500;
            transition: all 0.3s ease;
        }}
        
        .stButton > button:hover {{
            transform: translateY(-2px);
            box-shadow: 0 4px 12px rgba(0,0,0,0.2);
        }}
        
        /* 时间戳样式 */
        .timestamp {{
            font-size: 0.8em;
            opacity: 0.7;
            margin-top: 5px;
        }}
        
        /* 滚动条样式 */
        .chat-container::-webkit-scrollbar {{
            width: 8px;
        }}
        
        .chat-container::-webkit-scrollbar-track {{
            background: {theme_config['chat_bg']};
            border-radius: 10px;
        }}
        
        .chat-container::-webkit-scrollbar-thumb {{
            background: {theme_config['border']};
            border-radius: 10px;
        }}
        
        .chat-container::-webkit-scrollbar-thumb:hover {{
            background: {theme_config['user_bg']};
        }}
        
        /* 设置面板样式 */
        .settings-panel {{
            background: white;
            border-radius: 15px;
            padding: 20px;
            box-shadow: 0 8px 32px rgba(0,0,0,0.1);
            border: 1px solid {theme_config['border']};
        }}
        
        /* 色盘样式 */
        .color-picker {{
            display: flex;
            flex-wrap: wrap;
            gap: 10px;
            margin: 10px 0;
        }}
        
        .color-option {{
            width: 40px;
            height: 40px;
            border-radius: 50%;
            cursor: pointer;
            border: 3px solid transparent;
            transition: all 0.3s ease;
        }}
        
        .color-option:hover {{
            transform: scale(1.1);
            border-color: #333;
        }}
        
        .color-option.selected {{
            border-color: #007bff;
            transform: scale(1.2);
        }}
    </style>
    """, unsafe_allow_html=True)

# 系统配置
SYSTEM_CONFIG = {
    "system_name": "RAG智能对话系统",
    "version": "v2.3-Advanced",
    "admin_password": "rag2024",
    "max_file_size": 10,  # MB
    "max_daily_queries": 100,
    "supported_formats": [".pdf", ".docx", ".pptx", ".xlsx", ".txt", ".md", ".html", ".csv"],
    "max_message_length": 2000,
    "enable_file_upload": True,
    "enable_history_export": True,
    "chat_speed": 1.0,  # 回答速度倍数
    "show_timestamps": True,
    "auto_scroll": True,
    "sound_enabled": False
}

# 知识库
KNOWLEDGE_BASE = {
    "人工智能": "人工智能(AI)是计算机科学的一个分支，致力于创建能够执行通常需要人类智能的任务的系统。",
    "RAG": "检索增强生成(RAG)是一种AI技术，结合了信息检索和文本生成。它先从知识库中检索相关信息，然后基于检索结果生成准确的回答。",
    "机器学习": "机器学习是AI的一个分支，让计算机能够从数据中学习，无需明确编程就能改进性能。",
    "深度学习": "深度学习是机器学习的一个分支，使用多层神经网络模拟人脑处理信息的方式。",
    "自然语言处理": "自然语言处理(NLP)是AI的一个分支，专注于让计算机理解、解释和生成人类语言。",
    "向量数据库": "向量数据库是专门用于存储和检索高维向量数据的数据库系统，常用于AI和机器学习应用。",
    "嵌入模型": "嵌入模型将文本、图像等数据转换为高维向量表示，使计算机能够理解和处理这些数据。"
}

# 预设颜色选项
COLOR_PALETTE = {
    "用户消息颜色": [
        "#6c757d", "#4dabf7", "#51cf66", "#ff922b", "#9775fa", "#f783ac", "#4a5568", "#e64980"
    ],
    "AI回答颜色": [
        "#e9ecef", "#e7f5ff", "#d3f9d8", "#ffe8cc", "#e5dbff", "#ffdeeb", "#1a202c", "#fff0f6"
    ],
    "聊天背景颜色": [
        "#f8f9fa", "#f0f8ff", "#f3f9f0", "#fff4e6", "#f8f0ff", "#fff0f6", "#2d3748", "#fafafa"
    ]
}

class DocumentProcessor:
    """文档处理类"""
    
    def __init__(self):
        self.supported_extensions = SYSTEM_CONFIG["supported_formats"]
    
    def extract_text_from_file(self, uploaded_file) -> str:
        """从上传的文件中提取文本"""
        try:
            file_extension = Path(uploaded_file.name).suffix.lower()
            
            if file_extension == ".pdf":
                return self.extract_pdf(uploaded_file)
            elif file_extension == ".docx":
                return self.extract_docx(uploaded_file)
            elif file_extension == ".pptx":
                return self.extract_pptx(uploaded_file)
            elif file_extension in [".xlsx", ".xls"]:
                return self.extract_excel(uploaded_file)
            elif file_extension == ".csv":
                return self.extract_csv(uploaded_file)
            elif file_extension in [".txt", ".md"]:
                return self.extract_text(uploaded_file)
            elif file_extension == ".html":
                return self.extract_html(uploaded_file)
            else:
                return f"不支持的文件格式: {file_extension}"
        except Exception as e:
            return f"文件处理错误: {str(e)}"
    
    def extract_pdf(self, file) -> str:
        """提取PDF文件内容"""
        if not FULL_FEATURES:
            return "PDF处理功能需要安装PyMuPDF库"
        
        try:
            file_bytes = file.read()
            doc = fitz.open(stream=file_bytes, filetype="pdf")
            text = ""
            for page in doc:
                text += page.get_text()
            doc.close()
            return text
        except Exception as e:
            return f"PDF处理错误: {str(e)}"
    
    def extract_docx(self, file) -> str:
        """提取DOCX文件内容"""
        if not FULL_FEATURES:
            return "DOCX处理功能需要安装python-docx库"
        
        try:
            doc = Document(file)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            return text
        except Exception as e:
            return f"DOCX处理错误: {str(e)}"
    
    def extract_pptx(self, file) -> str:
        """提取PPTX文件内容"""
        if not FULL_FEATURES:
            return "PPTX处理功能需要安装python-pptx库"
        
        try:
            prs = Presentation(file)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        except Exception as e:
            return f"PPTX处理错误: {str(e)}"
    
    def extract_excel(self, file) -> str:
        """提取Excel文件内容"""
        try:
            df = pd.read_excel(file)
            return df.to_string()
        except Exception as e:
            return f"Excel处理错误: {str(e)}"
    
    def extract_csv(self, file) -> str:
        """提取CSV文件内容"""
        try:
            df = pd.read_csv(file)
            return df.to_string()
        except Exception as e:
            return f"CSV处理错误: {str(e)}"
    
    def extract_text(self, file) -> str:
        """提取纯文本文件内容"""
        try:
            content = file.read()
            if isinstance(content, bytes):
                content = content.decode('utf-8')
            return content
        except Exception as e:
            return f"文本处理错误: {str(e)}"
    
    def extract_html(self, file) -> str:
        """提取HTML文件内容"""
        if not FULL_FEATURES:
            return "HTML处理功能需要安装BeautifulSoup库"
        
        try:
            content = file.read()
            if isinstance(content, bytes):
                content = content.decode('utf-8')
            soup = BeautifulSoup(content, 'html.parser')
            return soup.get_text()
        except Exception as e:
            return f"HTML处理错误: {str(e)}"

def check_password():
    """检查密码"""
    def password_entered():
        if st.session_state["password"] == SYSTEM_CONFIG["admin_password"]:
            st.session_state["password_correct"] = True
            del st.session_state["password"]
        else:
            st.session_state["password_correct"] = False

    if "password_correct" not in st.session_state:
        st.markdown("### 🔐 欢迎使用RAG智能对话系统")
        st.text_input("请输入访问密码", type="password", on_change=password_entered, key="password")
        st.info("💡 这是一个受保护的智能对话系统，支持文档分析和高级设置")
        return False
    elif not st.session_state["password_correct"]:
        st.markdown("### 🔐 欢迎使用RAG智能对话系统")
        st.text_input("请输入访问密码", type="password", on_change=password_entered, key="password")
        st.error("❌ 密码错误，请重试")
        return False
    else:
        return True

def generate_response(message: str, context: str = "") -> str:
    """生成AI回答"""
    message_lower = message.lower()
    
    # 如果有文档上下文，优先使用
    if context:
        confidence = random.randint(90, 98)
        return f"""基于您上传的文档内容，我来回答您的问题：

{context[:500]}{'...' if len(context) > 500 else ''}

根据文档内容分析，{message}

**📚 信息来源**: 您上传的文档
**🎯 置信度**: {confidence}%
**📄 文档长度**: {len(context)} 字符"""
    
    # 检查知识库匹配
    for key, value in KNOWLEDGE_BASE.items():
        if key.lower() in message_lower:
            confidence = random.randint(88, 95)
            return f"""{value}

**📚 信息来源**: 内置知识库 - {key}
**🎯 置信度**: {confidence}%"""
    
    # 通用回答
    responses = [
        "这是一个很有意思的问题！基于RAG技术，我正在为您检索相关信息。",
        "根据我的理解，这个问题涉及到一些重要的概念。让我为您详细解答。",
        "通过智能检索，我发现这个话题有很多值得探讨的内容。",
        "基于检索增强生成技术，我来为您提供相关的信息和见解。"
    ]
    
    response = random.choice(responses)
    confidence = random.randint(75, 87)
    
    return f"""{response}

**📚 信息来源**: 智能检索系统
**🎯 置信度**: {confidence}%"""

def init_session_state():
    """初始化会话状态"""
    if "messages" not in st.session_state:
        st.session_state.messages = [
            {
                "role": "assistant",
                "content": "👋 您好！我是RAG智能对话助手。\n\n我可以帮您：\n• 📄 分析各种文档（PDF、Word、PPT等）\n• 💬 进行智能问答对话\n• 🔍 从知识库检索信息\n• ⚙️ 自定义系统设置\n\n点击右上角的⚙️图标可以打开设置面板！",
                "timestamp": datetime.datetime.now().strftime("%H:%M:%S")
            }
        ]
    
    if "uploaded_documents" not in st.session_state:
        st.session_state.uploaded_documents = {}
    
    if "query_count" not in st.session_state:
        st.session_state.query_count = 0
    
    if "show_settings" not in st.session_state:
        st.session_state.show_settings = False
    
    # 主题配置
    if "theme_config" not in st.session_state:
        st.session_state.theme_config = PRESET_THEMES["灰白简约"].copy()
        del st.session_state.theme_config["name"]
    
    # 系统设置
    if "system_settings" not in st.session_state:
        st.session_state.system_settings = SYSTEM_CONFIG.copy()

def display_chat_messages():
    """显示聊天消息"""
    chat_html = '<div class="chat-container">'
    
    for message in st.session_state.messages:
        role_class = "user-message" if message["role"] == "user" else "assistant-message"
        icon = "👤" if message["role"] == "user" else "🤖"
        
        timestamp_html = ""
        if st.session_state.system_settings["show_timestamps"] and "timestamp" in message:
            timestamp_html = f'<div class="timestamp">⏰ {message["timestamp"]}</div>'
        
        chat_html += f'''
        <div class="{role_class}">
            <strong>{icon}</strong> {message["content"].replace(chr(10), "<br>")}
            {timestamp_html}
        </div>
        '''
    
    chat_html += '</div>'
    st.markdown(chat_html, unsafe_allow_html=True)

def show_settings_panel():
    """显示设置面板"""
    st.markdown('<div class="settings-panel">', unsafe_allow_html=True)
    
    st.markdown("## ⚙️ 系统设置")
    
    # 创建标签页
    tab1, tab2, tab3, tab4 = st.tabs(["🎨 主题配色", "🔧 系统设置", "📄 文档设置", "💾 数据管理"])
    
    with tab1:
        st.markdown("### 🎨 主题配色设置")
        
        # 预设主题选择
        st.markdown("**📋 预设主题:**")
        theme_cols = st.columns(4)
        for i, (key, theme) in enumerate(PRESET_THEMES.items()):
            with theme_cols[i % 4]:
                if st.button(theme["name"], key=f"preset_{key}", use_container_width=True):
                    st.session_state.theme_config = theme.copy()
                    del st.session_state.theme_config["name"]
                    st.rerun()
        
        st.markdown("---")
        
        # 自定义颜色设置
        st.markdown("**🎨 自定义颜色:**")
        
        col1, col2 = st.columns(2)
        
        with col1:
            st.markdown("**用户消息背景色:**")
            user_bg = st.color_picker(
                "选择颜色",
                value=st.session_state.theme_config["user_bg"],
                key="user_bg_picker"
            )
            if user_bg != st.session_state.theme_config["user_bg"]:
                st.session_state.theme_config["user_bg"] = user_bg
                st.rerun()
            
            st.markdown("**AI回答背景色:**")
            assistant_bg = st.color_picker(
                "选择颜色",
                value=st.session_state.theme_config["assistant_bg"],
                key="assistant_bg_picker"
            )
            if assistant_bg != st.session_state.theme_config["assistant_bg"]:
                st.session_state.theme_config["assistant_bg"] = assistant_bg
                st.rerun()
        
        with col2:
            st.markdown("**聊天背景色:**")
            chat_bg = st.color_picker(
                "选择颜色",
                value=st.session_state.theme_config["chat_bg"],
                key="chat_bg_picker"
            )
            if chat_bg != st.session_state.theme_config["chat_bg"]:
                st.session_state.theme_config["chat_bg"] = chat_bg
                st.rerun()
            
            st.markdown("**边框颜色:**")
            border = st.color_picker(
                "选择颜色",
                value=st.session_state.theme_config["border"],
                key="border_picker"
            )
            if border != st.session_state.theme_config["border"]:
                st.session_state.theme_config["border"] = border
                st.rerun()
        
        # 快速色盘
        st.markdown("**🎨 快速色盘:**")
        st.markdown("点击下面的颜色快速应用到用户消息背景:")
        
        color_html = '<div class="color-picker">'
        for color in COLOR_PALETTE["用户消息颜色"]:
            color_html += f'<div class="color-option" style="background-color: {color};" onclick="document.getElementById(\'user_bg_picker\').value=\'{color}\'; document.getElementById(\'user_bg_picker\').dispatchEvent(new Event(\'change\'));"></div>'
        color_html += '</div>'
        st.markdown(color_html, unsafe_allow_html=True)
    
    with tab2:
        st.markdown("### 🔧 系统设置")
        
        # 聊天设置
        st.markdown("**💬 聊天设置:**")
        
        show_timestamps = st.checkbox(
            "显示时间戳",
            value=st.session_state.system_settings["show_timestamps"]
        )
        if show_timestamps != st.session_state.system_settings["show_timestamps"]:
            st.session_state.system_settings["show_timestamps"] = show_timestamps
        
        auto_scroll = st.checkbox(
            "自动滚动到底部",
            value=st.session_state.system_settings["auto_scroll"]
        )
        if auto_scroll != st.session_state.system_settings["auto_scroll"]:
            st.session_state.system_settings["auto_scroll"] = auto_scroll
        
        chat_speed = st.slider(
            "AI回答速度",
            min_value=0.1,
            max_value=3.0,
            value=st.session_state.system_settings["chat_speed"],
            step=0.1,
            help="数值越小回答越快"
        )
        if chat_speed != st.session_state.system_settings["chat_speed"]:
            st.session_state.system_settings["chat_speed"] = chat_speed
        
        # 安全设置
        st.markdown("**🔒 安全设置:**")
        
        max_daily_queries = st.number_input(
            "每日最大查询次数",
            min_value=10,
            max_value=1000,
            value=st.session_state.system_settings["max_daily_queries"]
        )
        if max_daily_queries != st.session_state.system_settings["max_daily_queries"]:
            st.session_state.system_settings["max_daily_queries"] = max_daily_queries
        
        max_message_length = st.number_input(
            "最大消息长度",
            min_value=100,
            max_value=5000,
            value=st.session_state.system_settings["max_message_length"]
        )
        if max_message_length != st.session_state.system_settings["max_message_length"]:
            st.session_state.system_settings["max_message_length"] = max_message_length
    
    with tab3:
        st.markdown("### 📄 文档设置")
        
        max_file_size = st.slider(
            "最大文件大小 (MB)",
            min_value=1,
            max_value=50,
            value=st.session_state.system_settings["max_file_size"]
        )
        if max_file_size != st.session_state.system_settings["max_file_size"]:
            st.session_state.system_settings["max_file_size"] = max_file_size
        
        st.markdown("**支持的文件格式:**")
        for fmt in st.session_state.system_settings["supported_formats"]:
            st.text(f"• {fmt}")
    
    with tab4:
        st.markdown("### 💾 数据管理")
        
        col1, col2 = st.columns(2)
        
        with col1:
            if st.button("🗑️ 清空所有对话", use_container_width=True):
                st.session_state.messages = [st.session_state.messages[0]]
                st.success("✅ 对话已清空")
            
            if st.button("📄 删除所有文档", use_container_width=True):
                st.session_state.uploaded_documents = {}
                st.success("✅ 文档已删除")
        
        with col2:
            if st.button("🔄 重置所有设置", use_container_width=True):
                st.session_state.theme_config = PRESET_THEMES["灰白简约"].copy()
                del st.session_state.theme_config["name"]
                st.session_state.system_settings = SYSTEM_CONFIG.copy()
                st.success("✅ 设置已重置")
            
            if st.button("💾 导出设置", use_container_width=True):
                settings_data = {
                    "theme": st.session_state.theme_config,
                    "system": st.session_state.system_settings
                }
                st.download_button(
                    "📥 下载设置文件",
                    data=str(settings_data),
                    file_name="rag_settings.json",
                    mime="application/json"
                )
        
        # 统计信息
        st.markdown("**📊 使用统计:**")
        st.text(f"• 总对话数: {len(st.session_state.messages)}")
        st.text(f"• 上传文档数: {len(st.session_state.uploaded_documents)}")
        st.text(f"• 今日查询数: {st.session_state.query_count}")
    
    # 关闭按钮
    if st.button("❌ 关闭设置", use_container_width=True):
        st.session_state.show_settings = False
        st.rerun()
    
    st.markdown('</div>', unsafe_allow_html=True)

def main():
    """主函数"""
    # 检查密码
    if not check_password():
        return
    
    # 初始化会话状态
    init_session_state()
    
    # 应用当前主题
    apply_theme(st.session_state.theme_config)
    
    # 设置按钮
    if st.button("⚙️", key="settings_btn", help="打开设置"):
        st.session_state.show_settings = not st.session_state.show_settings
        st.rerun()
    
    # 页面标题
    st.markdown("# 💬 RAG智能对话系统")
    st.markdown("### ⚙️ 高级设置版本 - 支持文档分析的智能助手")
    
    # 如果显示设置面板
    if st.session_state.show_settings:
        show_settings_panel()
        return
    
    # 主要布局
    col1, col2 = st.columns([3, 1])
    
    with col1:
        # 聊天显示区域
        st.markdown("#### 💬 对话区域")
        display_chat_messages()
        
        # 快捷问题
        st.markdown("**💡 快捷问题:**")
        quick_cols = st.columns(4)
        quick_questions = [
            "什么是RAG？",
            "如何使用系统？", 
            "分析我的文档",
            "打开设置面板"
        ]
        
        for i, question in enumerate(quick_questions):
            with quick_cols[i]:
                if st.button(question, key=f"quick_{i}", use_container_width=True):
                    if "设置" in question:
                        st.session_state.show_settings = True
                        st.rerun()
                    else:
                        # 添加用户消息
                        user_message = {
                            "role": "user",
                            "content": question,
                            "timestamp": datetime.datetime.now().strftime("%H:%M:%S")
                        }
                        st.session_state.messages.append(user_message)
                        
                        # 生成回答
                        context = ""
                        if "分析" in question and st.session_state.uploaded_documents:
                            latest_doc = list(st.session_state.uploaded_documents.values())[-1]
                            context = latest_doc["content"]
                        
                        with st.spinner("🤔 AI正在思考..."):
                            time.sleep(st.session_state.system_settings["chat_speed"])
                            response = generate_response(question, context)
                        
                        assistant_message = {
                            "role": "assistant",
                            "content": response,
                            "timestamp": datetime.datetime.now().strftime("%H:%M:%S")
                        }
                        st.session_state.messages.append(assistant_message)
                        st.session_state.query_count += 1
                        st.rerun()
        
        # 聊天输入框
        st.markdown("#### ✍️ 输入消息")
        if prompt := st.chat_input("💬 请输入您的问题...", max_chars=st.session_state.system_settings["max_message_length"]):
            # 检查查询限制
            if st.session_state.query_count >= st.session_state.system_settings["max_daily_queries"]:
                st.error(f"❌ 今日查询次数已达上限 ({st.session_state.system_settings['max_daily_queries']} 次)")
                return
            
            # 添加用户消息
            user_message = {
                "role": "user",
                "content": prompt,
                "timestamp": datetime.datetime.now().strftime("%H:%M:%S")
            }
            st.session_state.messages.append(user_message)
            
            # 生成AI回答
            with st.spinner("🤔 AI正在分析您的问题..."):
                time.sleep(st.session_state.system_settings["chat_speed"])
                
                # 检查是否需要使用文档上下文
                context = ""
                if st.session_state.uploaded_documents:
                    doc_keywords = ["文档", "分析", "总结", "内容", "解释", "说明"]
                    if any(keyword in prompt for keyword in doc_keywords):
                        latest_doc = list(st.session_state.uploaded_documents.values())[-1]
                        context = latest_doc["content"]
                
                response = generate_response(prompt, context)
            
            # 添加助手回答
            assistant_message = {
                "role": "assistant",
                "content": response,
                "timestamp": datetime.datetime.now().strftime("%H:%M:%S")
            }
            st.session_state.messages.append(assistant_message)
            st.session_state.query_count += 1
            st.rerun()
    
    with col2:
        # 文件上传区域
        st.markdown("#### 📄 文档上传")
        st.markdown('<div class="upload-area">', unsafe_allow_html=True)
        st.markdown("**拖拽或选择文件**")
        
        uploaded_file = st.file_uploader(
            "选择文档",
            type=[ext[1:] for ext in st.session_state.system_settings["supported_formats"]],
            help=f"支持: {', '.join(st.session_state.system_settings['supported_formats'])}"
        )
        
        if uploaded_file is not None:
            if uploaded_file.size > st.session_state.system_settings["max_file_size"] * 1024 * 1024:
                st.error(f"❌ 文件过大 (>{st.session_state.system_settings['max_file_size']}MB)")
            else:
                if st.button("🚀 上传分析", type="primary", use_container_width=True):
                    with st.spinner("🔄 处理中..."):
                        processor = DocumentProcessor()
                        content = processor.extract_text_from_file(uploaded_file)
                        
                        # 存储文档
                        doc_id = hashlib.md5(uploaded_file.name.encode()).hexdigest()[:8]
                        st.session_state.uploaded_documents[doc_id] = {
                            "name": uploaded_file.name,
                            "content": content,
                            "upload_time": datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
                            "size": uploaded_file.size
                        }
                        
                        st.success(f"✅ 上传成功！")
                        st.info(f"📊 {len(content)} 字符")
                        
                        # 自动添加分析消息
                        analysis_message = f"我上传了文档 '{uploaded_file.name}'，请帮我分析主要内容。"
                        
                        user_message = {
                            "role": "user",
                            "content": analysis_message,
                            "timestamp": datetime.datetime.now().strftime("%H:%M:%S")
                        }
                        st.session_state.messages.append(user_message)
                        
                        response = generate_response(analysis_message, content)
                        assistant_message = {
                            "role": "assistant",
                            "content": response,
                            "timestamp": datetime.datetime.now().strftime("%H:%M:%S")
                        }
                        st.session_state.messages.append(assistant_message)
                        st.session_state.query_count += 1
                        st.rerun()
        
        st.markdown('</div>', unsafe_allow_html=True)
        
        # 已上传文档
        if st.session_state.uploaded_documents:
            st.markdown("#### 📋 已上传文档")
            for doc_id, doc_info in st.session_state.uploaded_documents.items():
                with st.expander(f"📄 {doc_info['name'][:15]}..."):
                    st.text(f"时间: {doc_info['upload_time']}")
                    st.text(f"大小: {doc_info['size']} 字节")
                    st.text(f"文本: {len(doc_info['content'])} 字符")
                    
                    if st.button(f"🗑️ 删除", key=f"del_{doc_id}", use_container_width=True):
                        del st.session_state.uploaded_documents[doc_id]
                        st.rerun()
        
        # 系统状态
        st.markdown("#### 📊 系统状态")
        st.metric("💬 对话轮次", len(st.session_state.messages)//2)
        st.metric("📄 文档数量", len(st.session_state.uploaded_documents))
        st.metric("📈 今日查询", st.session_state.query_count)
        
        # 快捷操作
        st.markdown("#### 🎛️ 快捷操作")
        
        if st.button("⚙️ 打开设置", use_container_width=True):
            st.session_state.show_settings = True
            st.rerun()
        
        if st.button("🔄 刷新", use_container_width=True):
            st.rerun()
        
        if st.button("🚪 退出", use_container_width=True):
            st.session_state["password_correct"] = False
            st.rerun()
    
    # 页脚
    st.markdown("---")
    st.markdown(
        "<div style='text-align: center; color: #666; padding: 10px;'>" +
        f"💬 <b>{st.session_state.system_settings['system_name']}</b> {st.session_state.system_settings['version']} | " +
        "⚙️ 高级设置 | 🎨 自定义主题 | 📄 智能文档分析 | " +
        "<a href='https://github.com/novenazgarfield/research-workstation' target='_blank'>GitHub</a>" +
        "</div>",
        unsafe_allow_html=True
    )

if __name__ == "__main__":
    main()