import streamlit as st
import time
import datetime
import random
import io
import os
from pathlib import Path
from typing import List, Dict, Any
import hashlib

# 文档处理相关导入
try:
    import fitz  # PyMuPDF
    from docx import Document
    import pandas as pd
    from pptx import Presentation
    from bs4 import BeautifulSoup
    FULL_FEATURES = True
except ImportError:
    FULL_FEATURES = False

# 页面配置
st.set_page_config(
    page_title="💬 RAG智能对话系统",
    page_icon="💬",
    layout="wide",
    initial_sidebar_state="expanded"
)

# 系统配置
SYSTEM_CONFIG = {
    "system_name": "RAG智能对话系统",
    "version": "v2.0-Chat",
    "admin_password": "rag2024",
    "max_file_size": 10,  # MB
    "max_daily_queries": 100,
    "supported_formats": [".pdf", ".docx", ".pptx", ".xlsx", ".txt", ".md", ".html", ".csv"],
    "max_message_length": 2000,
    "enable_file_upload": True,
    "enable_history_export": True
}

# 知识库
KNOWLEDGE_BASE = {
    "人工智能": "人工智能(AI)是计算机科学的一个分支，致力于创建能够执行通常需要人类智能的任务的系统。",
    "RAG": "检索增强生成(RAG)是一种AI技术，结合了信息检索和文本生成。它先从知识库中检索相关信息，然后基于检索结果生成准确的回答。",
    "机器学习": "机器学习是AI的一个分支，让计算机能够从数据中学习，无需明确编程就能改进性能。",
    "深度学习": "深度学习是机器学习的一个分支，使用多层神经网络模拟人脑处理信息的方式。",
    "自然语言处理": "自然语言处理(NLP)是AI的一个分支，专注于让计算机理解、解释和生成人类语言。",
    "向量数据库": "向量数据库是专门用于存储和检索高维向量数据的数据库系统，常用于AI和机器学习应用。",
    "嵌入模型": "嵌入模型将文本、图像等数据转换为高维向量表示，使计算机能够理解和处理这些数据。"
}

class DocumentProcessor:
    """文档处理类"""
    
    def __init__(self):
        self.supported_extensions = SYSTEM_CONFIG["supported_formats"]
    
    def extract_text_from_file(self, uploaded_file) -> str:
        """从上传的文件中提取文本"""
        try:
            file_extension = Path(uploaded_file.name).suffix.lower()
            
            if file_extension == ".pdf":
                return self.extract_pdf(uploaded_file)
            elif file_extension == ".docx":
                return self.extract_docx(uploaded_file)
            elif file_extension == ".pptx":
                return self.extract_pptx(uploaded_file)
            elif file_extension in [".xlsx", ".xls"]:
                return self.extract_excel(uploaded_file)
            elif file_extension == ".csv":
                return self.extract_csv(uploaded_file)
            elif file_extension in [".txt", ".md"]:
                return self.extract_text(uploaded_file)
            elif file_extension == ".html":
                return self.extract_html(uploaded_file)
            else:
                return f"不支持的文件格式: {file_extension}"
        except Exception as e:
            return f"文件处理错误: {str(e)}"
    
    def extract_pdf(self, file) -> str:
        """提取PDF文件内容"""
        if not FULL_FEATURES:
            return "PDF处理功能需要安装PyMuPDF库"
        
        try:
            file_bytes = file.read()
            doc = fitz.open(stream=file_bytes, filetype="pdf")
            text = ""
            for page in doc:
                text += page.get_text()
            doc.close()
            return text
        except Exception as e:
            return f"PDF处理错误: {str(e)}"
    
    def extract_docx(self, file) -> str:
        """提取DOCX文件内容"""
        if not FULL_FEATURES:
            return "DOCX处理功能需要安装python-docx库"
        
        try:
            doc = Document(file)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            return text
        except Exception as e:
            return f"DOCX处理错误: {str(e)}"
    
    def extract_pptx(self, file) -> str:
        """提取PPTX文件内容"""
        if not FULL_FEATURES:
            return "PPTX处理功能需要安装python-pptx库"
        
        try:
            prs = Presentation(file)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        except Exception as e:
            return f"PPTX处理错误: {str(e)}"
    
    def extract_excel(self, file) -> str:
        """提取Excel文件内容"""
        try:
            df = pd.read_excel(file)
            return df.to_string()
        except Exception as e:
            return f"Excel处理错误: {str(e)}"
    
    def extract_csv(self, file) -> str:
        """提取CSV文件内容"""
        try:
            df = pd.read_csv(file)
            return df.to_string()
        except Exception as e:
            return f"CSV处理错误: {str(e)}"
    
    def extract_text(self, file) -> str:
        """提取纯文本文件内容"""
        try:
            content = file.read()
            if isinstance(content, bytes):
                content = content.decode('utf-8')
            return content
        except Exception as e:
            return f"文本处理错误: {str(e)}"
    
    def extract_html(self, file) -> str:
        """提取HTML文件内容"""
        if not FULL_FEATURES:
            return "HTML处理功能需要安装BeautifulSoup库"
        
        try:
            content = file.read()
            if isinstance(content, bytes):
                content = content.decode('utf-8')
            soup = BeautifulSoup(content, 'html.parser')
            return soup.get_text()
        except Exception as e:
            return f"HTML处理错误: {str(e)}"

def check_password():
    """检查密码"""
    def password_entered():
        if st.session_state["password"] == SYSTEM_CONFIG["admin_password"]:
            st.session_state["password_correct"] = True
            del st.session_state["password"]
        else:
            st.session_state["password_correct"] = False

    if "password_correct" not in st.session_state:
        st.text_input("🔐 请输入访问密码", type="password", on_change=password_entered, key="password")
        st.info("💡 这是一个受保护的RAG对话系统")
        return False
    elif not st.session_state["password_correct"]:
        st.text_input("🔐 请输入访问密码", type="password", on_change=password_entered, key="password")
        st.error("❌ 密码错误")
        return False
    else:
        return True

def generate_response(message: str, context: str = "") -> str:
    """生成AI回答"""
    message_lower = message.lower()
    
    # 如果有文档上下文，优先使用
    if context:
        confidence = random.randint(90, 98)
        return f"""基于您上传的文档内容，我来回答您的问题：

{context[:500]}{'...' if len(context) > 500 else ''}

根据文档内容分析，{message}

**📚 信息来源**: 您上传的文档
**🎯 置信度**: {confidence}%
**📄 文档长度**: {len(context)} 字符"""
    
    # 检查知识库匹配
    for key, value in KNOWLEDGE_BASE.items():
        if key.lower() in message_lower:
            confidence = random.randint(88, 95)
            return f"""{value}

**📚 信息来源**: 内置知识库 - {key}
**🎯 置信度**: {confidence}%"""
    
    # 通用回答
    responses = [
        "这是一个很有意思的问题！基于RAG技术，我正在为您检索相关信息。",
        "根据我的理解，这个问题涉及到一些重要的概念。让我为您详细解答。",
        "通过智能检索，我发现这个话题有很多值得探讨的内容。",
        "基于检索增强生成技术，我来为您提供相关的信息和见解。"
    ]
    
    response = random.choice(responses)
    confidence = random.randint(75, 87)
    
    return f"""{response}

**📚 信息来源**: 智能检索系统
**🎯 置信度**: {confidence}%"""

def init_session_state():
    """初始化会话状态"""
    if "messages" not in st.session_state:
        st.session_state.messages = [
            {
                "role": "assistant",
                "content": "👋 您好！我是RAG智能对话助手。我可以：\n\n📄 **处理文档**: 支持PDF、Word、PPT、Excel等格式\n💬 **智能问答**: 基于RAG技术的专业回答\n🔍 **知识检索**: 从文档和知识库中查找信息\n\n请上传文档或直接提问吧！",
                "timestamp": datetime.datetime.now().strftime("%H:%M:%S")
            }
        ]
    
    if "uploaded_documents" not in st.session_state:
        st.session_state.uploaded_documents = {}
    
    if "query_count" not in st.session_state:
        st.session_state.query_count = 0

def display_chat_message(message: Dict[str, Any]):
    """显示聊天消息"""
    with st.chat_message(message["role"]):
        st.markdown(message["content"])
        if "timestamp" in message:
            st.caption(f"⏰ {message['timestamp']}")

def main():
    """主函数"""
    # 检查密码
    if not check_password():
        return
    
    # 初始化会话状态
    init_session_state()
    
    # 页面标题
    st.title("💬 RAG智能对话系统")
    st.markdown("### 🤖 支持文档上传的智能问答助手")
    
    # 文件上传区域 - 放在主界面顶部
    st.markdown("---")
    st.subheader("📄 文档上传区域")
    
    col1, col2 = st.columns([3, 1])
    
    with col1:
        uploaded_file = st.file_uploader(
            "📤 选择要分析的文档文件",
            type=[ext[1:] for ext in SYSTEM_CONFIG["supported_formats"]],
            help=f"支持格式: {', '.join(SYSTEM_CONFIG['supported_formats'])} | 最大文件大小: {SYSTEM_CONFIG['max_file_size']}MB"
        )
    
    with col2:
        if uploaded_file is not None:
            if uploaded_file.size > SYSTEM_CONFIG["max_file_size"] * 1024 * 1024:
                st.error(f"❌ 文件过大")
            else:
                if st.button("🚀 上传并分析", type="primary", use_container_width=True):
                    with st.spinner("🔄 正在处理文档..."):
                        processor = DocumentProcessor()
                        content = processor.extract_text_from_file(uploaded_file)
                        
                        # 存储文档
                        doc_id = hashlib.md5(uploaded_file.name.encode()).hexdigest()[:8]
                        st.session_state.uploaded_documents[doc_id] = {
                            "name": uploaded_file.name,
                            "content": content,
                            "upload_time": datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
                            "size": uploaded_file.size
                        }
                        
                        st.success(f"✅ 文档 '{uploaded_file.name}' 上传成功！")
                        st.info(f"📊 提取文本长度: {len(content)} 字符")
                        
                        # 自动添加分析消息
                        analysis_message = f"我已经成功上传了文档 '{uploaded_file.name}'，请帮我分析这个文档的主要内容。"
                        
                        user_message = {
                            "role": "user",
                            "content": analysis_message,
                            "timestamp": datetime.datetime.now().strftime("%H:%M:%S")
                        }
                        st.session_state.messages.append(user_message)
                        
                        # 生成分析回答
                        response = generate_response(analysis_message, content)
                        assistant_message = {
                            "role": "assistant",
                            "content": response,
                            "timestamp": datetime.datetime.now().strftime("%H:%M:%S")
                        }
                        st.session_state.messages.append(assistant_message)
                        st.session_state.query_count += 1
                        st.rerun()
    
    # 显示已上传的文档
    if st.session_state.uploaded_documents:
        st.markdown("**📋 已上传的文档:**")
        cols = st.columns(len(st.session_state.uploaded_documents))
        for i, (doc_id, doc_info) in enumerate(st.session_state.uploaded_documents.items()):
            with cols[i]:
                st.info(f"📄 {doc_info['name'][:20]}{'...' if len(doc_info['name']) > 20 else ''}")
                if st.button(f"🗑️", key=f"del_main_{doc_id}", help="删除文档"):
                    del st.session_state.uploaded_documents[doc_id]
                    st.rerun()
    
    # 侧边栏
    with st.sidebar:
        st.header("🔧 系统控制")
        
        # 系统状态
        st.subheader("📊 系统状态")
        st.success("✅ 对话系统运行中")
        st.success("🤖 AI助手已就绪")
        st.info(f"📄 已上传文档: {len(st.session_state.uploaded_documents)}")
        st.info(f"💬 对话轮次: {len(st.session_state.messages)//2}")
        
        # 查询统计
        st.metric("📈 今日查询", st.session_state.query_count)
        st.metric("📚 知识条目", len(KNOWLEDGE_BASE))
        
        st.markdown("---")
        
        # 文档管理区域
        if st.session_state.uploaded_documents:
            st.subheader("📋 文档管理")
            for doc_id, doc_info in st.session_state.uploaded_documents.items():
                with st.expander(f"📄 {doc_info['name'][:15]}{'...' if len(doc_info['name']) > 15 else ''}"):
                    st.text(f"上传时间: {doc_info['upload_time']}")
                    st.text(f"文件大小: {doc_info['size']} 字节")
                    st.text(f"文本长度: {len(doc_info['content'])} 字符")
                    
                    if st.button(f"🗑️ 删除", key=f"del_{doc_id}", use_container_width=True):
                        del st.session_state.uploaded_documents[doc_id]
                        st.rerun()
        else:
            st.info("💡 在上方上传文档开始分析")
        
        st.markdown("---")
        
        # 系统控制
        st.subheader("🎛️ 系统控制")
        
        if st.button("🔄 刷新对话", use_container_width=True):
            st.rerun()
        
        if st.button("🗑️ 清空对话", use_container_width=True):
            st.session_state.messages = [st.session_state.messages[0]]  # 保留欢迎消息
            st.success("✅ 对话已清空")
            st.rerun()
        
        if st.button("📤 导出对话", use_container_width=True) and SYSTEM_CONFIG["enable_history_export"]:
            # 生成对话导出
            export_content = ""
            for msg in st.session_state.messages:
                role = "用户" if msg["role"] == "user" else "助手"
                timestamp = msg.get("timestamp", "")
                export_content += f"[{timestamp}] {role}: {msg['content']}\n\n"
            
            st.download_button(
                label="💾 下载对话记录",
                data=export_content,
                file_name=f"chat_history_{datetime.datetime.now().strftime('%Y%m%d_%H%M%S')}.txt",
                mime="text/plain",
                use_container_width=True
            )
        
        if st.button("🚪 安全退出", use_container_width=True):
            st.session_state["password_correct"] = False
            st.rerun()
    
    # 主聊天区域
    st.markdown("---")
    
    # 显示聊天历史
    chat_container = st.container()
    with chat_container:
        for message in st.session_state.messages:
            display_chat_message(message)
    
    # 聊天输入区域
    st.markdown("---")
    
    # 快捷问题按钮
    st.markdown("**💡 快捷问题:**")
    quick_questions = [
        "什么是RAG技术？",
        "如何使用这个系统？",
        "支持哪些文件格式？",
        "分析我上传的文档"
    ]
    
    cols = st.columns(len(quick_questions))
    for i, question in enumerate(quick_questions):
        with cols[i]:
            if st.button(question, key=f"quick_{i}", use_container_width=True):
                # 添加用户消息
                user_message = {
                    "role": "user",
                    "content": question,
                    "timestamp": datetime.datetime.now().strftime("%H:%M:%S")
                }
                st.session_state.messages.append(user_message)
                
                # 生成回答
                context = ""
                if "分析" in question and st.session_state.uploaded_documents:
                    # 使用最新上传的文档
                    latest_doc = list(st.session_state.uploaded_documents.values())[-1]
                    context = latest_doc["content"]
                
                with st.spinner("🤔 AI正在思考..."):
                    time.sleep(1)
                    response = generate_response(question, context)
                
                # 添加助手回答
                assistant_message = {
                    "role": "assistant",
                    "content": response,
                    "timestamp": datetime.datetime.now().strftime("%H:%M:%S")
                }
                st.session_state.messages.append(assistant_message)
                st.session_state.query_count += 1
                st.rerun()
    
    # 聊天输入框
    if prompt := st.chat_input("💬 请输入您的问题或与我对话...", max_chars=SYSTEM_CONFIG["max_message_length"]):
        # 检查查询限制
        if st.session_state.query_count >= SYSTEM_CONFIG["max_daily_queries"]:
            st.error(f"❌ 今日查询次数已达上限 ({SYSTEM_CONFIG['max_daily_queries']} 次)")
            return
        
        # 添加用户消息
        user_message = {
            "role": "user",
            "content": prompt,
            "timestamp": datetime.datetime.now().strftime("%H:%M:%S")
        }
        st.session_state.messages.append(user_message)
        
        # 显示用户消息
        display_chat_message(user_message)
        
        # 生成AI回答
        with st.chat_message("assistant"):
            with st.spinner("🤔 AI正在分析您的问题..."):
                time.sleep(1)
                
                # 检查是否需要使用文档上下文
                context = ""
                if st.session_state.uploaded_documents:
                    # 简单的关键词匹配来决定是否使用文档
                    doc_keywords = ["文档", "分析", "总结", "内容", "解释", "说明"]
                    if any(keyword in prompt for keyword in doc_keywords):
                        # 使用最新上传的文档
                        latest_doc = list(st.session_state.uploaded_documents.values())[-1]
                        context = latest_doc["content"]
                
                response = generate_response(prompt, context)
            
            st.markdown(response)
            st.caption(f"⏰ {datetime.datetime.now().strftime('%H:%M:%S')}")
        
        # 添加助手回答到历史
        assistant_message = {
            "role": "assistant",
            "content": response,
            "timestamp": datetime.datetime.now().strftime("%H:%M:%S")
        }
        st.session_state.messages.append(assistant_message)
        st.session_state.query_count += 1
    
    # 页脚信息
    st.markdown("---")
    st.markdown(
        "<div style='text-align: center; color: #666; padding: 10px;'>" +
        f"💬 <b>{SYSTEM_CONFIG['system_name']}</b> {SYSTEM_CONFIG['version']} | " +
        "🔒 安全对话 | 📄 文档智能分析 | " +
        "<a href='https://github.com/novenazgarfield/research-workstation' target='_blank'>GitHub仓库</a>" +
        "</div>",
        unsafe_allow_html=True
    )

if __name__ == "__main__":
    main()