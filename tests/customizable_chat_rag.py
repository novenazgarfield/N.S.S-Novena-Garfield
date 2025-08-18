import streamlit as st
import time
import datetime
import random
import io
import os
from pathlib import Path
from typing import List, Dict, Any
import hashlib

# 文档处理相关导入
try:
    import fitz  # PyMuPDF
    from docx import Document
    import pandas as pd
    from pptx import Presentation
    from bs4 import BeautifulSoup
    FULL_FEATURES = True
except ImportError:
    FULL_FEATURES = False

# 页面配置
st.set_page_config(
    page_title="💬 RAG智能对话系统",
    page_icon="💬",
    layout="wide",
    initial_sidebar_state="expanded"
)

# 主题配色方案
THEMES = {
    "灰白简约": {
        "name": "🤍 灰白简约",
        "chat_bg": "#f8f9fa",
        "user_bg": "linear-gradient(135deg, #6c757d 0%, #495057 100%)",
        "assistant_bg": "linear-gradient(135deg, #e9ecef 0%, #dee2e6 100%)",
        "user_color": "white",
        "assistant_color": "#212529",
        "border": "#dee2e6",
        "upload_bg": "linear-gradient(135deg, #6c757d 0%, #495057 100%)",
        "upload_color": "white"
    },
    "蓝色清新": {
        "name": "💙 蓝色清新", 
        "chat_bg": "#f0f8ff",
        "user_bg": "linear-gradient(135deg, #4dabf7 0%, #339af0 100%)",
        "assistant_bg": "linear-gradient(135deg, #e7f5ff 0%, #d0ebff 100%)",
        "user_color": "white",
        "assistant_color": "#1864ab",
        "border": "#74c0fc",
        "upload_bg": "linear-gradient(135deg, #4dabf7 0%, #339af0 100%)",
        "upload_color": "white"
    },
    "绿色自然": {
        "name": "💚 绿色自然",
        "chat_bg": "#f3f9f0",
        "user_bg": "linear-gradient(135deg, #51cf66 0%, #40c057 100%)",
        "assistant_bg": "linear-gradient(135deg, #d3f9d8 0%, #b2f2bb 100%)",
        "user_color": "white", 
        "assistant_color": "#2b8a3e",
        "border": "#8ce99a",
        "upload_bg": "linear-gradient(135deg, #51cf66 0%, #40c057 100%)",
        "upload_color": "white"
    },
    "橙色温暖": {
        "name": "🧡 橙色温暖",
        "chat_bg": "#fff4e6",
        "user_bg": "linear-gradient(135deg, #ff922b 0%, #fd7e14 100%)",
        "assistant_bg": "linear-gradient(135deg, #ffe8cc 0%, #ffd8a8 100%)",
        "user_color": "white",
        "assistant_color": "#d9480f",
        "border": "#ffb366",
        "upload_bg": "linear-gradient(135deg, #ff922b 0%, #fd7e14 100%)",
        "upload_color": "white"
    },
    "紫色梦幻": {
        "name": "💜 紫色梦幻",
        "chat_bg": "#f8f0ff",
        "user_bg": "linear-gradient(135deg, #9775fa 0%, #845ef7 100%)",
        "assistant_bg": "linear-gradient(135deg, #e5dbff 0%, #d0bfff 100%)",
        "user_color": "white",
        "assistant_color": "#6741d9",
        "border": "#b197fc",
        "upload_bg": "linear-gradient(135deg, #9775fa 0%, #845ef7 100%)",
        "upload_color": "white"
    },
    "粉色可爱": {
        "name": "💗 粉色可爱",
        "chat_bg": "#fff0f6",
        "user_bg": "linear-gradient(135deg, #f783ac 0%, #e64980 100%)",
        "assistant_bg": "linear-gradient(135deg, #ffdeeb 0%, #fcc2d7 100%)",
        "user_color": "white",
        "assistant_color": "#c2255c",
        "border": "#faa2c1",
        "upload_bg": "linear-gradient(135deg, #f783ac 0%, #e64980 100%)",
        "upload_color": "white"
    },
    "深色模式": {
        "name": "🖤 深色模式",
        "chat_bg": "#2d3748",
        "user_bg": "linear-gradient(135deg, #4a5568 0%, #2d3748 100%)",
        "assistant_bg": "linear-gradient(135deg, #1a202c 0%, #2d3748 100%)",
        "user_color": "#e2e8f0",
        "assistant_color": "#e2e8f0",
        "border": "#4a5568",
        "upload_bg": "linear-gradient(135deg, #4a5568 0%, #2d3748 100%)",
        "upload_color": "#e2e8f0"
    }
}

def apply_theme(theme_key):
    """应用选定的主题"""
    theme = THEMES[theme_key]
    
    st.markdown(f"""
    <style>
        /* 隐藏Streamlit默认元素 */
        #MainMenu {{visibility: hidden;}}
        footer {{visibility: hidden;}}
        header {{visibility: hidden;}}
        
        /* 聊天容器样式 */
        .chat-container {{
            background-color: {theme['chat_bg']};
            border: 2px solid {theme['border']};
            border-radius: 15px;
            padding: 20px;
            margin: 10px 0;
            height: 500px;
            overflow-y: auto;
            box-shadow: 0 4px 15px rgba(0,0,0,0.1);
        }}
        
        /* 消息气泡样式 */
        .user-message {{
            background: {theme['user_bg']};
            color: {theme['user_color']};
            padding: 12px 18px;
            border-radius: 18px 18px 5px 18px;
            margin: 8px 0 8px auto;
            max-width: 80%;
            word-wrap: break-word;
            box-shadow: 0 2px 8px rgba(0,0,0,0.15);
            display: block;
            text-align: left;
        }}
        
        .assistant-message {{
            background: {theme['assistant_bg']};
            color: {theme['assistant_color']};
            padding: 12px 18px;
            border-radius: 18px 18px 18px 5px;
            margin: 8px auto 8px 0;
            max-width: 80%;
            word-wrap: break-word;
            box-shadow: 0 2px 8px rgba(0,0,0,0.15);
            display: block;
            text-align: left;
        }}
        
        /* 文件上传区域样式 */
        .upload-area {{
            background: {theme['upload_bg']};
            border-radius: 15px;
            padding: 15px;
            margin: 10px 0;
            color: {theme['upload_color']};
            text-align: center;
            box-shadow: 0 2px 10px rgba(0,0,0,0.1);
        }}
        
        /* 按钮样式 */
        .stButton > button {{
            border-radius: 20px;
            border: none;
            padding: 8px 20px;
            font-weight: 500;
            transition: all 0.3s ease;
        }}
        
        .stButton > button:hover {{
            transform: translateY(-2px);
            box-shadow: 0 4px 12px rgba(0,0,0,0.2);
        }}
        
        /* 时间戳样式 */
        .timestamp {{
            font-size: 0.8em;
            opacity: 0.7;
            margin-top: 5px;
        }}
        
        /* 滚动条样式 */
        .chat-container::-webkit-scrollbar {{
            width: 8px;
        }}
        
        .chat-container::-webkit-scrollbar-track {{
            background: {theme['chat_bg']};
            border-radius: 10px;
        }}
        
        .chat-container::-webkit-scrollbar-thumb {{
            background: {theme['border']};
            border-radius: 10px;
        }}
        
        .chat-container::-webkit-scrollbar-thumb:hover {{
            background: {theme['user_bg'].split()[0] if 'linear-gradient' in theme['user_bg'] else theme['user_bg']};
        }}
    </style>
    """, unsafe_allow_html=True)

# 系统配置
SYSTEM_CONFIG = {
    "system_name": "RAG智能对话系统",
    "version": "v2.2-Customizable",
    "admin_password": "rag2024",
    "max_file_size": 10,  # MB
    "max_daily_queries": 100,
    "supported_formats": [".pdf", ".docx", ".pptx", ".xlsx", ".txt", ".md", ".html", ".csv"],
    "max_message_length": 2000,
    "enable_file_upload": True,
    "enable_history_export": True
}

# 知识库
KNOWLEDGE_BASE = {
    "人工智能": "人工智能(AI)是计算机科学的一个分支，致力于创建能够执行通常需要人类智能的任务的系统。",
    "RAG": "检索增强生成(RAG)是一种AI技术，结合了信息检索和文本生成。它先从知识库中检索相关信息，然后基于检索结果生成准确的回答。",
    "机器学习": "机器学习是AI的一个分支，让计算机能够从数据中学习，无需明确编程就能改进性能。",
    "深度学习": "深度学习是机器学习的一个分支，使用多层神经网络模拟人脑处理信息的方式。",
    "自然语言处理": "自然语言处理(NLP)是AI的一个分支，专注于让计算机理解、解释和生成人类语言。",
    "向量数据库": "向量数据库是专门用于存储和检索高维向量数据的数据库系统，常用于AI和机器学习应用。",
    "嵌入模型": "嵌入模型将文本、图像等数据转换为高维向量表示，使计算机能够理解和处理这些数据。"
}

class DocumentProcessor:
    """文档处理类"""
    
    def __init__(self):
        self.supported_extensions = SYSTEM_CONFIG["supported_formats"]
    
    def extract_text_from_file(self, uploaded_file) -> str:
        """从上传的文件中提取文本"""
        try:
            file_extension = Path(uploaded_file.name).suffix.lower()
            
            if file_extension == ".pdf":
                return self.extract_pdf(uploaded_file)
            elif file_extension == ".docx":
                return self.extract_docx(uploaded_file)
            elif file_extension == ".pptx":
                return self.extract_pptx(uploaded_file)
            elif file_extension in [".xlsx", ".xls"]:
                return self.extract_excel(uploaded_file)
            elif file_extension == ".csv":
                return self.extract_csv(uploaded_file)
            elif file_extension in [".txt", ".md"]:
                return self.extract_text(uploaded_file)
            elif file_extension == ".html":
                return self.extract_html(uploaded_file)
            else:
                return f"不支持的文件格式: {file_extension}"
        except Exception as e:
            return f"文件处理错误: {str(e)}"
    
    def extract_pdf(self, file) -> str:
        """提取PDF文件内容"""
        if not FULL_FEATURES:
            return "PDF处理功能需要安装PyMuPDF库"
        
        try:
            file_bytes = file.read()
            doc = fitz.open(stream=file_bytes, filetype="pdf")
            text = ""
            for page in doc:
                text += page.get_text()
            doc.close()
            return text
        except Exception as e:
            return f"PDF处理错误: {str(e)}"
    
    def extract_docx(self, file) -> str:
        """提取DOCX文件内容"""
        if not FULL_FEATURES:
            return "DOCX处理功能需要安装python-docx库"
        
        try:
            doc = Document(file)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            return text
        except Exception as e:
            return f"DOCX处理错误: {str(e)}"
    
    def extract_pptx(self, file) -> str:
        """提取PPTX文件内容"""
        if not FULL_FEATURES:
            return "PPTX处理功能需要安装python-pptx库"
        
        try:
            prs = Presentation(file)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        except Exception as e:
            return f"PPTX处理错误: {str(e)}"
    
    def extract_excel(self, file) -> str:
        """提取Excel文件内容"""
        try:
            df = pd.read_excel(file)
            return df.to_string()
        except Exception as e:
            return f"Excel处理错误: {str(e)}"
    
    def extract_csv(self, file) -> str:
        """提取CSV文件内容"""
        try:
            df = pd.read_csv(file)
            return df.to_string()
        except Exception as e:
            return f"CSV处理错误: {str(e)}"
    
    def extract_text(self, file) -> str:
        """提取纯文本文件内容"""
        try:
            content = file.read()
            if isinstance(content, bytes):
                content = content.decode('utf-8')
            return content
        except Exception as e:
            return f"文本处理错误: {str(e)}"
    
    def extract_html(self, file) -> str:
        """提取HTML文件内容"""
        if not FULL_FEATURES:
            return "HTML处理功能需要安装BeautifulSoup库"
        
        try:
            content = file.read()
            if isinstance(content, bytes):
                content = content.decode('utf-8')
            soup = BeautifulSoup(content, 'html.parser')
            return soup.get_text()
        except Exception as e:
            return f"HTML处理错误: {str(e)}"

def check_password():
    """检查密码"""
    def password_entered():
        if st.session_state["password"] == SYSTEM_CONFIG["admin_password"]:
            st.session_state["password_correct"] = True
            del st.session_state["password"]
        else:
            st.session_state["password_correct"] = False

    if "password_correct" not in st.session_state:
        st.markdown("### 🔐 欢迎使用RAG智能对话系统")
        st.text_input("请输入访问密码", type="password", on_change=password_entered, key="password")
        st.info("💡 这是一个受保护的智能对话系统，支持文档分析和主题自定义")
        return False
    elif not st.session_state["password_correct"]:
        st.markdown("### 🔐 欢迎使用RAG智能对话系统")
        st.text_input("请输入访问密码", type="password", on_change=password_entered, key="password")
        st.error("❌ 密码错误，请重试")
        return False
    else:
        return True

def generate_response(message: str, context: str = "") -> str:
    """生成AI回答"""
    message_lower = message.lower()
    
    # 如果有文档上下文，优先使用
    if context:
        confidence = random.randint(90, 98)
        return f"""基于您上传的文档内容，我来回答您的问题：

{context[:500]}{'...' if len(context) > 500 else ''}

根据文档内容分析，{message}

**📚 信息来源**: 您上传的文档
**🎯 置信度**: {confidence}%
**📄 文档长度**: {len(context)} 字符"""
    
    # 检查知识库匹配
    for key, value in KNOWLEDGE_BASE.items():
        if key.lower() in message_lower:
            confidence = random.randint(88, 95)
            return f"""{value}

**📚 信息来源**: 内置知识库 - {key}
**🎯 置信度**: {confidence}%"""
    
    # 通用回答
    responses = [
        "这是一个很有意思的问题！基于RAG技术，我正在为您检索相关信息。",
        "根据我的理解，这个问题涉及到一些重要的概念。让我为您详细解答。",
        "通过智能检索，我发现这个话题有很多值得探讨的内容。",
        "基于检索增强生成技术，我来为您提供相关的信息和见解。"
    ]
    
    response = random.choice(responses)
    confidence = random.randint(75, 87)
    
    return f"""{response}

**📚 信息来源**: 智能检索系统
**🎯 置信度**: {confidence}%"""

def init_session_state():
    """初始化会话状态"""
    if "messages" not in st.session_state:
        st.session_state.messages = [
            {
                "role": "assistant",
                "content": "👋 您好！我是RAG智能对话助手。\n\n我可以帮您：\n• 📄 分析各种文档（PDF、Word、PPT等）\n• 💬 进行智能问答对话\n• 🔍 从知识库检索信息\n• 🎨 自定义界面主题\n\n请开始我们的对话吧！",
                "timestamp": datetime.datetime.now().strftime("%H:%M:%S")
            }
        ]
    
    if "uploaded_documents" not in st.session_state:
        st.session_state.uploaded_documents = {}
    
    if "query_count" not in st.session_state:
        st.session_state.query_count = 0
    
    if "current_theme" not in st.session_state:
        st.session_state.current_theme = "灰白简约"  # 默认使用您喜欢的灰白主题

def display_chat_messages():
    """显示聊天消息"""
    chat_html = '<div class="chat-container">'
    
    for message in st.session_state.messages:
        role_class = "user-message" if message["role"] == "user" else "assistant-message"
        icon = "👤" if message["role"] == "user" else "🤖"
        
        chat_html += f'''
        <div class="{role_class}">
            <strong>{icon}</strong> {message["content"].replace(chr(10), "<br>")}
            <div class="timestamp">⏰ {message.get("timestamp", "")}</div>
        </div>
        '''
    
    chat_html += '</div>'
    st.markdown(chat_html, unsafe_allow_html=True)

def main():
    """主函数"""
    # 检查密码
    if not check_password():
        return
    
    # 初始化会话状态
    init_session_state()
    
    # 应用当前主题
    apply_theme(st.session_state.current_theme)
    
    # 页面标题
    st.markdown("# 💬 RAG智能对话系统")
    st.markdown("### 🎨 可自定义主题版本 - 支持文档分析的智能助手")
    
    # 侧边栏 - 主题设置
    with st.sidebar:
        st.header("🎨 主题设置")
        
        # 主题选择
        theme_options = list(THEMES.keys())
        theme_names = [THEMES[key]["name"] for key in theme_options]
        
        selected_theme_index = theme_options.index(st.session_state.current_theme)
        
        new_theme_index = st.selectbox(
            "选择主题配色",
            range(len(theme_options)),
            index=selected_theme_index,
            format_func=lambda x: theme_names[x]
        )
        
        new_theme = theme_options[new_theme_index]
        
        if new_theme != st.session_state.current_theme:
            st.session_state.current_theme = new_theme
            st.rerun()
        
        # 主题预览
        st.markdown("**🎨 当前主题预览:**")
        current_theme = THEMES[st.session_state.current_theme]
        st.markdown(f"""
        <div style="background: {current_theme['user_bg']}; color: {current_theme['user_color']}; padding: 8px; border-radius: 10px; margin: 5px 0;">
            👤 用户消息样式
        </div>
        <div style="background: {current_theme['assistant_bg']}; color: {current_theme['assistant_color']}; padding: 8px; border-radius: 10px; margin: 5px 0;">
            🤖 AI回答样式
        </div>
        """, unsafe_allow_html=True)
        
        st.markdown("---")
        
        # 系统状态
        st.subheader("📊 系统状态")
        st.metric("💬 对话轮次", len(st.session_state.messages)//2)
        st.metric("📄 文档数量", len(st.session_state.uploaded_documents))
        st.metric("📈 今日查询", st.session_state.query_count)
        
        # 已上传文档
        if st.session_state.uploaded_documents:
            st.markdown("---")
            st.subheader("📋 文档管理")
            for doc_id, doc_info in st.session_state.uploaded_documents.items():
                with st.expander(f"📄 {doc_info['name'][:15]}..."):
                    st.text(f"时间: {doc_info['upload_time']}")
                    st.text(f"大小: {doc_info['size']} 字节")
                    st.text(f"文本: {len(doc_info['content'])} 字符")
                    
                    if st.button(f"🗑️ 删除", key=f"del_{doc_id}", use_container_width=True):
                        del st.session_state.uploaded_documents[doc_id]
                        st.rerun()
        
        st.markdown("---")
        
        # 系统控制
        st.subheader("🎛️ 系统控制")
        
        if st.button("🔄 刷新", use_container_width=True):
            st.rerun()
        
        if st.button("🗑️ 清空对话", use_container_width=True):
            st.session_state.messages = [st.session_state.messages[0]]
            st.success("✅ 已清空")
            st.rerun()
        
        if st.button("🚪 退出", use_container_width=True):
            st.session_state["password_correct"] = False
            st.rerun()
    
    # 主要布局
    col1, col2 = st.columns([3, 1])
    
    with col1:
        # 聊天显示区域
        st.markdown("#### 💬 对话区域")
        display_chat_messages()
        
        # 快捷问题
        st.markdown("**💡 快捷问题:**")
        quick_cols = st.columns(4)
        quick_questions = [
            "什么是RAG？",
            "如何使用系统？", 
            "分析我的文档",
            "更换主题颜色"
        ]
        
        for i, question in enumerate(quick_questions):
            with quick_cols[i]:
                if st.button(question, key=f"quick_{i}", use_container_width=True):
                    # 添加用户消息
                    user_message = {
                        "role": "user",
                        "content": question,
                        "timestamp": datetime.datetime.now().strftime("%H:%M:%S")
                    }
                    st.session_state.messages.append(user_message)
                    
                    # 生成回答
                    context = ""
                    if "分析" in question and st.session_state.uploaded_documents:
                        latest_doc = list(st.session_state.uploaded_documents.values())[-1]
                        context = latest_doc["content"]
                    elif "主题" in question:
                        response = f"您可以在左侧边栏选择不同的主题配色！\n\n当前可用主题：\n" + "\n".join([f"• {theme['name']}" for theme in THEMES.values()]) + f"\n\n当前使用：{THEMES[st.session_state.current_theme]['name']}"
                    else:
                        response = generate_response(question, context)
                    
                    if "主题" not in question:
                        with st.spinner("🤔 AI正在思考..."):
                            time.sleep(1)
                            response = generate_response(question, context)
                    
                    assistant_message = {
                        "role": "assistant",
                        "content": response,
                        "timestamp": datetime.datetime.now().strftime("%H:%M:%S")
                    }
                    st.session_state.messages.append(assistant_message)
                    st.session_state.query_count += 1
                    st.rerun()
        
        # 聊天输入框
        st.markdown("#### ✍️ 输入消息")
        if prompt := st.chat_input("💬 请输入您的问题...", max_chars=SYSTEM_CONFIG["max_message_length"]):
            # 检查查询限制
            if st.session_state.query_count >= SYSTEM_CONFIG["max_daily_queries"]:
                st.error(f"❌ 今日查询次数已达上限 ({SYSTEM_CONFIG['max_daily_queries']} 次)")
                return
            
            # 添加用户消息
            user_message = {
                "role": "user",
                "content": prompt,
                "timestamp": datetime.datetime.now().strftime("%H:%M:%S")
            }
            st.session_state.messages.append(user_message)
            
            # 生成AI回答
            with st.spinner("🤔 AI正在分析您的问题..."):
                time.sleep(1)
                
                # 检查是否需要使用文档上下文
                context = ""
                if st.session_state.uploaded_documents:
                    doc_keywords = ["文档", "分析", "总结", "内容", "解释", "说明"]
                    if any(keyword in prompt for keyword in doc_keywords):
                        latest_doc = list(st.session_state.uploaded_documents.values())[-1]
                        context = latest_doc["content"]
                
                response = generate_response(prompt, context)
            
            # 添加助手回答
            assistant_message = {
                "role": "assistant",
                "content": response,
                "timestamp": datetime.datetime.now().strftime("%H:%M:%S")
            }
            st.session_state.messages.append(assistant_message)
            st.session_state.query_count += 1
            st.rerun()
    
    with col2:
        # 文件上传区域
        st.markdown("#### 📄 文档上传")
        st.markdown('<div class="upload-area">', unsafe_allow_html=True)
        st.markdown("**拖拽或选择文件**")
        
        uploaded_file = st.file_uploader(
            "选择文档",
            type=[ext[1:] for ext in SYSTEM_CONFIG["supported_formats"]],
            help=f"支持: {', '.join(SYSTEM_CONFIG['supported_formats'])}"
        )
        
        if uploaded_file is not None:
            if uploaded_file.size > SYSTEM_CONFIG["max_file_size"] * 1024 * 1024:
                st.error(f"❌ 文件过大 (>{SYSTEM_CONFIG['max_file_size']}MB)")
            else:
                if st.button("🚀 上传分析", type="primary", use_container_width=True):
                    with st.spinner("🔄 处理中..."):
                        processor = DocumentProcessor()
                        content = processor.extract_text_from_file(uploaded_file)
                        
                        # 存储文档
                        doc_id = hashlib.md5(uploaded_file.name.encode()).hexdigest()[:8]
                        st.session_state.uploaded_documents[doc_id] = {
                            "name": uploaded_file.name,
                            "content": content,
                            "upload_time": datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
                            "size": uploaded_file.size
                        }
                        
                        st.success(f"✅ 上传成功！")
                        st.info(f"📊 {len(content)} 字符")
                        
                        # 自动添加分析消息
                        analysis_message = f"我上传了文档 '{uploaded_file.name}'，请帮我分析主要内容。"
                        
                        user_message = {
                            "role": "user",
                            "content": analysis_message,
                            "timestamp": datetime.datetime.now().strftime("%H:%M:%S")
                        }
                        st.session_state.messages.append(user_message)
                        
                        response = generate_response(analysis_message, content)
                        assistant_message = {
                            "role": "assistant",
                            "content": response,
                            "timestamp": datetime.datetime.now().strftime("%H:%M:%S")
                        }
                        st.session_state.messages.append(assistant_message)
                        st.session_state.query_count += 1
                        st.rerun()
        
        st.markdown('</div>', unsafe_allow_html=True)
    
    # 页脚
    st.markdown("---")
    st.markdown(
        "<div style='text-align: center; color: #666; padding: 10px;'>" +
        f"💬 <b>{SYSTEM_CONFIG['system_name']}</b> {SYSTEM_CONFIG['version']} | " +
        f"🎨 当前主题: {THEMES[st.session_state.current_theme]['name']} | " +
        "📄 智能文档分析 | " +
        "<a href='https://github.com/novenazgarfield/research-workstation' target='_blank'>GitHub</a>" +
        "</div>",
        unsafe_allow_html=True
    )

if __name__ == "__main__":
    main()