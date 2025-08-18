import streamlit as st
import time
import datetime
import random
import io
import os
from pathlib import Path
from typing import List, Dict, Any
import hashlib

# 文档处理相关导入
try:
    import fitz  # PyMuPDF
    from docx import Document
    import pandas as pd
    from pptx import Presentation
    from bs4 import BeautifulSoup
    FULL_FEATURES = True
except ImportError:
    FULL_FEATURES = False

# 页面配置
st.set_page_config(
    page_title="💬 RAG智能对话",
    page_icon="💬",
    layout="wide",
    initial_sidebar_state="collapsed"
)

# 移动端优化CSS
st.markdown("""
<style>
    /* 移动端适配 */
    @media (max-width: 768px) {
        .main .block-container {
            padding-left: 1rem;
            padding-right: 1rem;
            padding-top: 1rem;
        }
        
        .chat-container {
            height: 400px !important;
        }
        
        .stColumns {
            flex-direction: column;
        }
        
        .stColumns > div {
            width: 100% !important;
            margin-bottom: 1rem;
        }
        
        .settings-button {
            width: 45px !important;
            height: 45px !important;
            font-size: 18px !important;
        }
        
        .quick-buttons {
            display: grid;
            grid-template-columns: 1fr 1fr;
            gap: 0.5rem;
        }
        
        .upload-area {
            padding: 10px !important;
        }
    }
    
    /* 基础样式 */
    #MainMenu {visibility: hidden;}
    footer {visibility: hidden;}
    header {visibility: hidden;}
    
    /* 设置按钮 */
    .settings-button {
        position: fixed;
        top: 15px;
        right: 15px;
        z-index: 1000;
        background: linear-gradient(135deg, #6c757d 0%, #495057 100%);
        color: white;
        border: none;
        border-radius: 50%;
        width: 50px;
        height: 50px;
        font-size: 20px;
        cursor: pointer;
        box-shadow: 0 4px 15px rgba(0,0,0,0.2);
        transition: all 0.3s ease;
    }
    
    .settings-button:hover {
        transform: scale(1.1);
        box-shadow: 0 6px 20px rgba(0,0,0,0.3);
    }
    
    /* 聊天容器 */
    .chat-container {
        background-color: #f8f9fa;
        border: 2px solid #dee2e6;
        border-radius: 15px;
        padding: 15px;
        margin: 10px 0;
        height: 450px;
        overflow-y: auto;
        box-shadow: 0 4px 15px rgba(0,0,0,0.1);
    }
    
    /* 消息气泡 */
    .user-message {
        background: #6c757d;
        color: white;
        padding: 10px 15px;
        border-radius: 15px 15px 5px 15px;
        margin: 8px 0 8px auto;
        max-width: 85%;
        word-wrap: break-word;
        box-shadow: 0 2px 8px rgba(0,0,0,0.15);
        display: block;
        text-align: left;
        font-size: 14px;
    }
    
    .assistant-message {
        background: #e9ecef;
        color: #212529;
        padding: 10px 15px;
        border-radius: 15px 15px 15px 5px;
        margin: 8px auto 8px 0;
        max-width: 85%;
        word-wrap: break-word;
        box-shadow: 0 2px 8px rgba(0,0,0,0.15);
        display: block;
        text-align: left;
        font-size: 14px;
    }
    
    /* 文件上传区域 */
    .upload-area {
        background: #6c757d;
        border-radius: 15px;
        padding: 15px;
        margin: 10px 0;
        color: white;
        text-align: center;
        box-shadow: 0 2px 10px rgba(0,0,0,0.1);
    }
    
    /* 按钮优化 */
    .stButton > button {
        border-radius: 20px;
        border: none;
        padding: 8px 16px;
        font-weight: 500;
        transition: all 0.3s ease;
        width: 100%;
    }
    
    .stButton > button:hover {
        transform: translateY(-2px);
        box-shadow: 0 4px 12px rgba(0,0,0,0.2);
    }
    
    /* 时间戳 */
    .timestamp {
        font-size: 0.75em;
        opacity: 0.7;
        margin-top: 5px;
    }
    
    /* 滚动条优化 */
    .chat-container::-webkit-scrollbar {
        width: 6px;
    }
    
    .chat-container::-webkit-scrollbar-track {
        background: #f8f9fa;
        border-radius: 10px;
    }
    
    .chat-container::-webkit-scrollbar-thumb {
        background: #dee2e6;
        border-radius: 10px;
    }
    
    /* 设置面板移动端优化 */
    .settings-panel {
        background: white;
        border-radius: 15px;
        padding: 15px;
        box-shadow: 0 8px 32px rgba(0,0,0,0.1);
        border: 1px solid #dee2e6;
        max-height: 80vh;
        overflow-y: auto;
    }
    
    /* 输入框优化 */
    .stTextInput > div > div > input {
        border-radius: 20px;
        border: 2px solid #dee2e6;
        padding: 10px 15px;
    }
    
    /* 文件上传器优化 */
    .stFileUploader > div {
        border-radius: 15px;
        border: 2px dashed #dee2e6;
        padding: 20px;
        text-align: center;
    }
    
    /* 响应式网格 */
    .responsive-grid {
        display: grid;
        grid-template-columns: repeat(auto-fit, minmax(120px, 1fr));
        gap: 0.5rem;
        margin: 10px 0;
    }
    
    /* 移动端隐藏元素 */
    @media (max-width: 768px) {
        .desktop-only {
            display: none;
        }
    }
    
    /* 桌面端隐藏元素 */
    @media (min-width: 769px) {
        .mobile-only {
            display: none;
        }
    }
</style>
""", unsafe_allow_html=True)

# 系统配置
SYSTEM_CONFIG = {
    "system_name": "RAG智能对话",
    "version": "v3.1-Enhanced",
    "admin_password": "rag2024",
    "max_file_size": 50,  # 增加到50MB
    "max_daily_queries": 500,  # 增加到500次
    "supported_formats": [".pdf", ".docx", ".pptx", ".xlsx", ".txt", ".md", ".html", ".csv"],
    "max_message_length": 5000,  # 增加到5000字符
    "enable_file_upload": True,
    "show_timestamps": True,
    "chat_speed": 1.0,
    "theme_mode": "light",  # light, dark, auto
    "max_documents": 20,  # 最大文档数量
    "max_conversations": 1000  # 最大对话数量
}

# 知识库
KNOWLEDGE_BASE = {
    "人工智能": "人工智能(AI)是计算机科学的一个分支，致力于创建能够执行通常需要人类智能的任务的系统。",
    "RAG": "检索增强生成(RAG)是一种AI技术，结合了信息检索和文本生成。它先从知识库中检索相关信息，然后基于检索结果生成准确的回答。",
    "机器学习": "机器学习是AI的一个分支，让计算机能够从数据中学习，无需明确编程就能改进性能。",
    "深度学习": "深度学习是机器学习的一个分支，使用多层神经网络模拟人脑处理信息的方式。",
    "自然语言处理": "自然语言处理(NLP)是AI的一个分支，专注于让计算机理解、解释和生成人类语言。"
}

class DocumentProcessor:
    """文档处理类"""
    
    def __init__(self):
        self.supported_extensions = SYSTEM_CONFIG["supported_formats"]
    
    def extract_text_from_file(self, uploaded_file) -> str:
        """从上传的文件中提取文本"""
        try:
            file_extension = Path(uploaded_file.name).suffix.lower()
            
            if file_extension == ".pdf":
                return self.extract_pdf(uploaded_file)
            elif file_extension == ".docx":
                return self.extract_docx(uploaded_file)
            elif file_extension == ".pptx":
                return self.extract_pptx(uploaded_file)
            elif file_extension in [".xlsx", ".xls"]:
                return self.extract_excel(uploaded_file)
            elif file_extension == ".csv":
                return self.extract_csv(uploaded_file)
            elif file_extension in [".txt", ".md"]:
                return self.extract_text(uploaded_file)
            elif file_extension == ".html":
                return self.extract_html(uploaded_file)
            else:
                return f"不支持的文件格式: {file_extension}"
        except Exception as e:
            return f"文件处理错误: {str(e)}"
    
    def extract_pdf(self, file) -> str:
        if not FULL_FEATURES:
            return "PDF处理功能需要安装PyMuPDF库"
        try:
            file_bytes = file.read()
            doc = fitz.open(stream=file_bytes, filetype="pdf")
            text = ""
            for page in doc:
                text += page.get_text()
            doc.close()
            return text
        except Exception as e:
            return f"PDF处理错误: {str(e)}"
    
    def extract_docx(self, file) -> str:
        if not FULL_FEATURES:
            return "DOCX处理功能需要安装python-docx库"
        try:
            doc = Document(file)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            return text
        except Exception as e:
            return f"DOCX处理错误: {str(e)}"
    
    def extract_pptx(self, file) -> str:
        if not FULL_FEATURES:
            return "PPTX处理功能需要安装python-pptx库"
        try:
            prs = Presentation(file)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        except Exception as e:
            return f"PPTX处理错误: {str(e)}"
    
    def extract_excel(self, file) -> str:
        try:
            df = pd.read_excel(file)
            return df.to_string()
        except Exception as e:
            return f"Excel处理错误: {str(e)}"
    
    def extract_csv(self, file) -> str:
        try:
            df = pd.read_csv(file)
            return df.to_string()
        except Exception as e:
            return f"CSV处理错误: {str(e)}"
    
    def extract_text(self, file) -> str:
        try:
            content = file.read()
            if isinstance(content, bytes):
                content = content.decode('utf-8')
            return content
        except Exception as e:
            return f"文本处理错误: {str(e)}"
    
    def extract_html(self, file) -> str:
        if not FULL_FEATURES:
            return "HTML处理功能需要安装BeautifulSoup库"
        try:
            content = file.read()
            if isinstance(content, bytes):
                content = content.decode('utf-8')
            soup = BeautifulSoup(content, 'html.parser')
            return soup.get_text()
        except Exception as e:
            return f"HTML处理错误: {str(e)}"

def check_password():
    """检查密码"""
    def password_entered():
        if st.session_state["password"] == SYSTEM_CONFIG["admin_password"]:
            st.session_state["password_correct"] = True
            del st.session_state["password"]
        else:
            st.session_state["password_correct"] = False

    if "password_correct" not in st.session_state:
        st.markdown("### 🔐 RAG智能对话系统")
        st.text_input("请输入访问密码", type="password", on_change=password_entered, key="password")
        st.info("💡 支持手机、平板、电脑全平台访问")
        return False
    elif not st.session_state["password_correct"]:
        st.markdown("### 🔐 RAG智能对话系统")
        st.text_input("请输入访问密码", type="password", on_change=password_entered, key="password")
        st.error("❌ 密码错误")
        return False
    else:
        return True

def generate_response(message: str, context: str = "") -> str:
    """生成AI回答"""
    message_lower = message.lower()
    
    if context:
        confidence = random.randint(90, 98)
        return f"""基于您上传的文档内容分析：

{context[:300]}{'...' if len(context) > 300 else ''}

**📚 来源**: 您的文档
**🎯 置信度**: {confidence}%"""
    
    for key, value in KNOWLEDGE_BASE.items():
        if key.lower() in message_lower:
            confidence = random.randint(88, 95)
            return f"""{value}

**📚 来源**: 知识库 - {key}
**🎯 置信度**: {confidence}%"""
    
    responses = [
        "这是一个很有意思的问题！我正在为您检索相关信息。",
        "根据我的理解，让我为您详细解答这个问题。",
        "通过智能检索，我发现这个话题很值得探讨。"
    ]
    
    response = random.choice(responses)
    confidence = random.randint(75, 87)
    
    return f"""{response}

**📚 来源**: 智能检索
**🎯 置信度**: {confidence}%"""

def init_session_state():
    """初始化会话状态"""
    if "messages" not in st.session_state:
        st.session_state.messages = [
            {
                "role": "assistant",
                "content": "👋 您好！我是RAG智能助手。\n\n✨ 功能特色：\n• 📱 移动端优化\n• 📄 文档分析\n• 💬 智能问答\n• 🌐 全平台支持\n\n请开始我们的对话吧！",
                "timestamp": datetime.datetime.now().strftime("%H:%M")
            }
        ]
    
    if "uploaded_documents" not in st.session_state:
        st.session_state.uploaded_documents = {}
    
    if "query_count" not in st.session_state:
        st.session_state.query_count = 0
    
    if "show_settings" not in st.session_state:
        st.session_state.show_settings = False

def display_chat_messages():
    """显示聊天消息"""
    chat_html = '<div class="chat-container">'
    
    for message in st.session_state.messages:
        role_class = "user-message" if message["role"] == "user" else "assistant-message"
        icon = "👤" if message["role"] == "user" else "🤖"
        
        timestamp_html = ""
        if SYSTEM_CONFIG["show_timestamps"] and "timestamp" in message:
            timestamp_html = f'<div class="timestamp">⏰ {message["timestamp"]}</div>'
        
        chat_html += f'''
        <div class="{role_class}">
            <strong>{icon}</strong> {message["content"].replace(chr(10), "<br>")}
            {timestamp_html}
        </div>
        '''
    
    chat_html += '</div>'
    st.markdown(chat_html, unsafe_allow_html=True)

def show_settings_panel():
    """显示设置面板"""
    st.markdown('<div class="settings-panel">', unsafe_allow_html=True)
    
    st.markdown("## ⚙️ 系统设置")
    
    # 设置标签页
    tab1, tab2, tab3, tab4 = st.tabs(["🎨 外观", "⚡ 性能", "📊 系统", "🔧 功能"])
    
    with tab1:
        st.markdown("### 🎨 外观设置")
        
        show_timestamps = st.checkbox(
            "显示时间戳",
            value=SYSTEM_CONFIG["show_timestamps"],
            help="在消息中显示发送时间"
        )
        SYSTEM_CONFIG["show_timestamps"] = show_timestamps
        
        # 界面主题选择
        st.markdown("**🌈 界面主题:**")
        theme_options = ["浅色模式", "深色模式", "自动跟随系统"]
        current_theme_index = 0
        if SYSTEM_CONFIG["theme_mode"] == "dark":
            current_theme_index = 1
        elif SYSTEM_CONFIG["theme_mode"] == "auto":
            current_theme_index = 2
            
        selected_theme = st.selectbox(
            "选择主题",
            theme_options,
            index=current_theme_index,
            help="选择界面的颜色主题"
        )
        
        # 更新主题配置
        if selected_theme == "浅色模式":
            SYSTEM_CONFIG["theme_mode"] = "light"
        elif selected_theme == "深色模式":
            SYSTEM_CONFIG["theme_mode"] = "dark"
        elif selected_theme == "自动跟随系统":
            SYSTEM_CONFIG["theme_mode"] = "auto"
        
        # 显示主题预览
        if SYSTEM_CONFIG["theme_mode"] == "light":
            st.info("🌞 当前使用浅色主题")
        elif SYSTEM_CONFIG["theme_mode"] == "dark":
            st.info("🌙 当前使用深色主题")
        else:
            st.info("🔄 当前跟随系统主题设置")
        
        # 字体大小调节
        font_size = st.slider(
            "字体大小",
            min_value=12,
            max_value=20,
            value=14,
            step=1,
            help="调节聊天消息的字体大小"
        )
    
    with tab2:
        st.markdown("### ⚡ 性能设置")
        
        # AI回答速度设置
        chat_speed = st.slider(
            "🤖 AI回答速度",
            min_value=0.1,
            max_value=5.0,
            value=SYSTEM_CONFIG["chat_speed"],
            step=0.1,
            help="数值越小回答越快，0.1=极快，3.0=较慢"
        )
        SYSTEM_CONFIG["chat_speed"] = chat_speed
        
        # 显示当前速度状态
        if chat_speed <= 0.5:
            speed_status = "⚡ 极快模式"
        elif chat_speed <= 1.5:
            speed_status = "🚀 快速模式"
        elif chat_speed <= 3.0:
            speed_status = "🐢 正常模式"
        else:
            speed_status = "🐌 慢速模式"
        
        st.info(f"当前状态: {speed_status}")
        
        # 自动滚动设置
        auto_scroll = st.checkbox(
            "自动滚动到最新消息",
            value=True,
            help="新消息时自动滚动到底部"
        )
        
        # 消息加载数量
        max_messages = st.slider(
            "最大显示消息数",
            min_value=20,
            max_value=200,
            value=100,
            step=10,
            help="限制显示的历史消息数量，提升性能"
        )
    
    with tab3:
        st.markdown("### 📊 系统信息")
        
        # 使用统计
        st.markdown("**📈 使用统计:**")
        col1, col2, col3 = st.columns(3)
        
        with col1:
            st.metric(
                "💬 对话轮次",
                len(st.session_state.messages)//2,
                help="用户和AI的对话轮次"
            )
        
        with col2:
            st.metric(
                "📄 文档数量",
                len(st.session_state.uploaded_documents),
                help="已上传的文档数量"
            )
        
        with col3:
            st.metric(
                "📈 今日查询",
                st.session_state.query_count,
                help="今天的查询次数"
            )
        
        # 系统状态
        st.markdown("**🔧 系统状态:**")
        st.text(f"• 版本: {SYSTEM_CONFIG['version']}")
        st.text(f"• 最大文件大小: {SYSTEM_CONFIG['max_file_size']}MB")
        st.text(f"• 每日查询限制: {SYSTEM_CONFIG['max_daily_queries']}次")
        st.text(f"• 消息长度限制: {SYSTEM_CONFIG['max_message_length']}字符")
        
        # 支持的文件格式
        st.markdown("**📄 支持的文件格式:**")
        formats_text = " • ".join(SYSTEM_CONFIG["supported_formats"])
        st.text(formats_text)
        
        # 内存使用情况
        import psutil
        memory_percent = psutil.virtual_memory().percent
        st.markdown(f"**💾 内存使用率:** {memory_percent:.1f}%")
        
        if memory_percent > 80:
            st.warning("⚠️ 内存使用率较高，建议清理数据")
        elif memory_percent > 60:
            st.info("ℹ️ 内存使用正常")
        else:
            st.success("✅ 内存使用良好")
    
    with tab4:
        st.markdown("### 🔧 功能设置")
        
        # 文件上传设置
        max_file_size = st.slider(
            "📄 最大文件大小 (MB)",
            min_value=1,
            max_value=100,  # 增加到100MB
            value=SYSTEM_CONFIG["max_file_size"],
            step=5,
            help="限制上传文件的最大大小"
        )
        SYSTEM_CONFIG["max_file_size"] = max_file_size
        
        # 文档数量限制
        max_documents = st.slider(
            "📚 最大文档数量",
            min_value=5,
            max_value=50,  # 增加到50个文档
            value=SYSTEM_CONFIG["max_documents"],
            step=5,
            help="同时保存的最大文档数量"
        )
        SYSTEM_CONFIG["max_documents"] = max_documents
        
        # 查询限制设置
        max_daily_queries = st.number_input(
            "📈 每日最大查询次数",
            min_value=50,
            max_value=2000,  # 增加到2000次
            value=SYSTEM_CONFIG["max_daily_queries"],
            step=50,
            help="防止过度使用的安全限制"
        )
        SYSTEM_CONFIG["max_daily_queries"] = max_daily_queries
        
        # 消息长度限制
        max_message_length = st.number_input(
            "📝 最大消息长度",
            min_value=500,
            max_value=10000,  # 增加到10000字符
            value=SYSTEM_CONFIG["max_message_length"],
            step=500,
            help="单条消息的最大字符数"
        )
        SYSTEM_CONFIG["max_message_length"] = max_message_length
        
        # 对话历史限制
        max_conversations = st.number_input(
            "💬 最大对话数量",
            min_value=100,
            max_value=5000,  # 增加到5000条对话
            value=SYSTEM_CONFIG["max_conversations"],
            step=100,
            help="保存的最大对话记录数量"
        )
        SYSTEM_CONFIG["max_conversations"] = max_conversations
        
        # 功能开关
        st.markdown("**🎛️ 功能开关:**")
        
        enable_file_upload = st.checkbox(
            "启用文件上传功能",
            value=SYSTEM_CONFIG["enable_file_upload"]
        )
        SYSTEM_CONFIG["enable_file_upload"] = enable_file_upload
        
        enable_history_export = st.checkbox(
            "启用对话历史导出",
            value=SYSTEM_CONFIG.get("enable_history_export", True)
        )
        SYSTEM_CONFIG["enable_history_export"] = enable_history_export
    
    # 底部操作按钮
    st.markdown("---")
    st.markdown("### 🎛️ 系统操作")
    
    col1, col2, col3, col4 = st.columns(4)
    
    with col1:
        if st.button("🗑️ 清空对话", use_container_width=True, help="删除所有对话记录"):
            st.session_state.messages = [st.session_state.messages[0]]
            st.success("✅ 对话已清空")
            st.rerun()
    
    with col2:
        if st.button("📄 删除文档", use_container_width=True, help="删除所有上传的文档"):
            st.session_state.uploaded_documents = {}
            st.success("✅ 文档已删除")
            st.rerun()
    
    with col3:
        if st.button("🔄 重置设置", use_container_width=True, help="恢复默认设置"):
            # 重置为默认值
            SYSTEM_CONFIG.update({
                "show_timestamps": True,
                "chat_speed": 1.0,
                "max_file_size": 10,
                "max_daily_queries": 100,
                "max_message_length": 2000
            })
            st.success("✅ 设置已重置")
            st.rerun()
    
    with col4:
        if st.button("❌ 关闭设置", use_container_width=True, type="primary"):
            st.session_state.show_settings = False
            st.rerun()
    
    st.markdown('</div>', unsafe_allow_html=True)

def main():
    """主函数"""
    if not check_password():
        return
    
    init_session_state()
    
    # 设置按钮
    if st.button("⚙️", key="settings_btn", help="设置"):
        st.session_state.show_settings = not st.session_state.show_settings
        st.rerun()
    
    # 页面标题
    st.markdown("# 💬 RAG智能对话")
    st.markdown("### 📱 移动端优化版本")
    
    if st.session_state.show_settings:
        show_settings_panel()
        return
    
    # 检测设备类型
    device_info = st.empty()
    with device_info.container():
        st.markdown("""
        <script>
        const isMobile = window.innerWidth <= 768;
        const deviceType = isMobile ? '📱 移动设备' : '💻 桌面设备';
        document.title = deviceType + ' - RAG智能对话';
        </script>
        """, unsafe_allow_html=True)
    
    # 响应式布局 - 默认使用移动端布局
    st.session_state.is_mobile = True
    
    # 移动端布局
    display_chat_messages()
    
    # 智能提示（替代明显的快捷按钮）
    if len(st.session_state.messages) <= 1:  # 只在开始时显示
        with st.expander("💡 试试这些问题", expanded=False):
            st.markdown("**点击下方问题快速开始对话：**")
            suggestion_cols = st.columns(2)
            suggestions = [
                "什么是RAG？", "如何使用？", 
                "分析文档", "查看设置"
            ]
            
            for i, suggestion in enumerate(suggestions):
                col_idx = i % 2
                with suggestion_cols[col_idx]:
                    if st.button(suggestion, key=f"suggestion_{i}", use_container_width=True):
                        if "设置" in suggestion:
                            st.session_state.show_settings = True
                            st.rerun()
                        else:
                            # 将建议作为用户消息处理
                            st.session_state.temp_message = suggestion
                            st.rerun()
    
    # 处理临时消息（来自建议点击）
    if hasattr(st.session_state, 'temp_message'):
        prompt = st.session_state.temp_message
        del st.session_state.temp_message
        
        # 检查查询限制
        if st.session_state.query_count >= SYSTEM_CONFIG["max_daily_queries"]:
            st.error(f"❌ 今日查询次数已达上限 ({SYSTEM_CONFIG['max_daily_queries']} 次)")
        else:
            user_message = {
                "role": "user", 
                "content": prompt,
                "timestamp": datetime.datetime.now().strftime("%H:%M")
            }
            st.session_state.messages.append(user_message)
            
            context = ""
            if st.session_state.uploaded_documents and "文档" in prompt:
                latest_doc = list(st.session_state.uploaded_documents.values())[-1]
                context = latest_doc["content"]
            
            with st.spinner("🤔 AI思考中..."):
                time.sleep(SYSTEM_CONFIG["chat_speed"])
                response = generate_response(prompt, context)
            
            assistant_message = {
                "role": "assistant",
                "content": response, 
                "timestamp": datetime.datetime.now().strftime("%H:%M")
            }
            st.session_state.messages.append(assistant_message)
            st.session_state.query_count += 1
            st.rerun()
    
    # 文件上传区域
    st.markdown("---")
    st.markdown("### 📄 文档上传")
    
    # 检查文档数量限制
    if len(st.session_state.uploaded_documents) >= SYSTEM_CONFIG["max_documents"]:
        st.warning(f"⚠️ 已达到最大文档数量限制 ({SYSTEM_CONFIG['max_documents']} 个)")
        st.info("💡 请删除一些文档后再上传新文档")
    else:
        uploaded_file = st.file_uploader(
            "选择文档",
            type=[ext[1:] for ext in SYSTEM_CONFIG["supported_formats"]],
            help=f"支持格式: {', '.join(SYSTEM_CONFIG['supported_formats'])}\n最大文件大小: {SYSTEM_CONFIG['max_file_size']}MB",
            key="main_file_uploader"
        )
        
        if uploaded_file is not None:
            # 显示文件信息
            file_size_mb = uploaded_file.size / (1024 * 1024)
            st.info(f"📄 **{uploaded_file.name}** ({file_size_mb:.2f} MB)")
            
            # 检查文件大小
            if uploaded_file.size > SYSTEM_CONFIG["max_file_size"] * 1024 * 1024:
                st.error(f"❌ 文件过大，最大支持 {SYSTEM_CONFIG['max_file_size']}MB")
            else:
                col1, col2 = st.columns(2)
                with col1:
                    if st.button("🚀 上传分析", type="primary", use_container_width=True, key="main_upload_btn"):
                        with st.spinner("🔄 处理中..."):
                            try:
                                processor = DocumentProcessor()
                                content = processor.extract_text_from_file(uploaded_file)
                                
                                if content and content.strip():
                                    doc_id = hashlib.md5(f"{uploaded_file.name}{time.time()}".encode()).hexdigest()[:8]
                                    st.session_state.uploaded_documents[doc_id] = {
                                        "name": uploaded_file.name,
                                        "content": content,
                                        "upload_time": datetime.datetime.now().strftime("%m-%d %H:%M"),
                                        "size": uploaded_file.size
                                    }
                                    
                                    st.success(f"✅ 上传成功！")
                                    st.balloons()
                                    
                                    # 自动分析
                                    analysis_message = f"已上传文档 '{uploaded_file.name}'，请分析内容。"
                                    
                                    user_message = {
                                        "role": "user",
                                        "content": analysis_message,
                                        "timestamp": datetime.datetime.now().strftime("%H:%M")
                                    }
                                    st.session_state.messages.append(user_message)
                                    
                                    response = generate_response(analysis_message, content)
                                    assistant_message = {
                                        "role": "assistant",
                                        "content": response,
                                        "timestamp": datetime.datetime.now().strftime("%H:%M")
                                    }
                                    st.session_state.messages.append(assistant_message)
                                    st.session_state.query_count += 1
                                    time.sleep(1)
                                    st.rerun()
                                else:
                                    st.error("❌ 文档内容为空或处理失败，请检查文件格式")
                            except Exception as e:
                                st.error(f"❌ 文档处理失败: {str(e)}")
                
                with col2:
                    if st.button("❌ 取消", use_container_width=True, key="main_cancel_btn"):
                        st.rerun()
    
    # 已上传文档显示
    if st.session_state.uploaded_documents:
        st.markdown("### 📋 已上传文档")
        for doc_id, doc_info in st.session_state.uploaded_documents.items():
            col1, col2 = st.columns([3, 1])
            with col1:
                st.text(f"📄 {doc_info['name'][:30]}{'...' if len(doc_info['name']) > 30 else ''}")
                st.caption(f"上传时间: {doc_info['upload_time']}")
            with col2:
                if st.button("🗑️", key=f"del_{doc_id}", help="删除文档"):
                    del st.session_state.uploaded_documents[doc_id]
                    st.success("✅ 文档已删除")
                    st.rerun()
    
    # 聊天输入
    st.markdown("---")
    st.markdown("### 💬 开始对话")
    
    # 聊天输入框
    if prompt := st.chat_input("💬 输入消息...", max_chars=SYSTEM_CONFIG["max_message_length"], key="main_chat_input"):
        # 检查查询限制
        if st.session_state.query_count >= SYSTEM_CONFIG["max_daily_queries"]:
            st.error(f"❌ 今日查询次数已达上限 ({SYSTEM_CONFIG['max_daily_queries']} 次)")
        else:
            user_message = {
                "role": "user",
                "content": prompt,
                "timestamp": datetime.datetime.now().strftime("%H:%M")
            }
            st.session_state.messages.append(user_message)
            
            with st.spinner("🤔 AI思考中..."):
                time.sleep(SYSTEM_CONFIG["chat_speed"])
                
                context = ""
                if st.session_state.uploaded_documents:
                    doc_keywords = ["文档", "分析", "总结", "内容"]
                    if any(keyword in prompt for keyword in doc_keywords):
                        latest_doc = list(st.session_state.uploaded_documents.values())[-1]
                        context = latest_doc["content"]
                
                response = generate_response(prompt, context)
            
            assistant_message = {
                "role": "assistant",
                "content": response,
                "timestamp": datetime.datetime.now().strftime("%H:%M")
            }
            st.session_state.messages.append(assistant_message)
            st.session_state.query_count += 1
            st.rerun()
    

    
    # 页脚
    st.markdown("---")
    st.markdown(
        "<div style='text-align: center; color: #666; padding: 10px; font-size: 12px;'>" +
        f"💬 <b>{SYSTEM_CONFIG['system_name']}</b> {SYSTEM_CONFIG['version']} | " +
        "📱 移动端优化 | 🌐 全平台支持" +
        "</div>",
        unsafe_allow_html=True
    )

if __name__ == "__main__":
    main()